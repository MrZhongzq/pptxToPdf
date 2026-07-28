"""生成自检用的内置 pptx。开发时跑一次，产物提交进仓库。

运行时不需要 python-pptx（它只在 requirements-dev.txt），所以这个
文件是资产而不是运行时生成的。

用法（工作目录 backend/）：
    .venv/Scripts/python.exe scripts/make_selftest_pptx.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

MARKER = "pptx2pdf selftest"
DEST = Path(__file__).resolve().parent.parent / "app" / "assets" / "selftest.pptx"


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    frame = box.text_frame
    frame.text = MARKER
    frame.paragraphs[0].runs[0].font.size = Pt(40)

    DEST.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(DEST))
    print(f"wrote {DEST} ({DEST.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
