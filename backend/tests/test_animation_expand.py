"""动画分步展开。

fixture 的 p:timing 结构逐字照抄自用户那份 59 页真实课件（slide29 的
整形状动画、slide42 的段落级动画）——五期栽过一次：合成 fixture 与真实
PowerPoint 结构不符，导致「剥离后能转」全是假绿，真实课件一上传就 406。
"""

import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.animation_expand import (
    MAX_STEPS_PER_SLIDE,
    P,
    Target,
    expand_animations,
    parse_steps,
)

# 真实结构：mainSeq > childTnLst > par(每个 par 一次点击) > ... >
# cTn{presetClass=entr, nodeType=clickEffect} > cBhvr > tgtEl > spTgt
_STEP_TEMPLATE = (
    '<p:par><p:cTn fill="hold">'
    '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
    "<p:childTnLst><p:par><p:cTn fill=\"hold\">"
    '<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
    "<p:childTnLst><p:par>"
    '<p:cTn presetID="21" presetClass="entr" fill="hold" nodeType="clickEffect">'
    "<p:childTnLst>{behaviors}</p:childTnLst>"
    "</p:cTn>"
    "</p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>"
    "</p:cTn></p:par>"
)


def _behavior(spid: str, para: int | None = None) -> str:
    tgt = f'<p:spTgt spid="{spid}"/>'
    if para is not None:
        tgt = (
            f'<p:spTgt spid="{spid}"><p:txEl>'
            f'<p:pRg st="{para}" end="{para}"/></p:txEl></p:spTgt>'
        )
    return f"<p:cBhvr><p:cTn/><p:tgtEl>{tgt}</p:tgtEl></p:cBhvr>"


def _timing(*steps: str) -> str:
    return (
        "<p:timing><p:tnLst><p:par>"
        '<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst>'
        '<p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f"<p:childTnLst>{''.join(steps)}</p:childTnLst>"
        "</p:cTn></p:seq>"
        "</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"
    )


def _deck(path: Path, pages: int = 2) -> Path:
    prs = Presentation()
    for i in range(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"PAGE-{i + 1}"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    prs.save(str(path))
    return path


def _inject(src: Path, dst: Path, slide_no: int, timing_xml: str, shapes_xml: str = "") -> Path:
    """往指定页塞入 timing 与额外形状。"""
    target = f"ppt/slides/slide{slide_no}.xml"
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            raw = zin.read(item.filename)
            if item.filename == target:
                text = raw.decode("utf-8")
                if shapes_xml:
                    text = text.replace("</p:spTree>", shapes_xml + "</p:spTree>")
                text = text.replace("</p:sld>", timing_xml + "</p:sld>")
                raw = text.encode("utf-8")
            zout.writestr(item, raw)
    return dst


def _shape(shape_id: str, name: str, paragraphs: int = 1) -> str:
    paras = "".join(
        f'<a:p><a:r><a:t>行{i}</a:t></a:r></a:p>' for i in range(paragraphs)
    )
    return (
        f"<p:sp><p:nvSpPr>"
        f'<p:cNvPr id="{shape_id}" name="{name}"/><p:cNvSpPr/><p:nvPr/>'
        f"</p:nvSpPr>"
        f'<p:spPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="100" cy="100"/></a:xfrm></p:spPr>'
        f"<p:txBody><a:bodyPr/><a:lstStyle/>{paras}</p:txBody>"
        f"</p:sp>"
    )


def _pages(path: Path) -> int:
    with zipfile.ZipFile(path) as zf:
        root = ET.fromstring(zf.read("ppt/presentation.xml"))
        return len(root.find(f"{P}sldIdLst").findall(f"{P}sldId"))


def _slide_shape_ids(path: Path, part: str) -> list[str]:
    with zipfile.ZipFile(path) as zf:
        root = ET.fromstring(zf.read(part))
    tree = root.find(f"{P}cSld/{P}spTree")
    out = []
    for sp in tree:
        c = sp.find(f".//{P}cNvPr")
        if c is not None:
            out.append(c.get("id"))
    return out


# ---- 解析 ----


def test_parse_steps_groups_by_click(tmp_path):
    xml = f"<p:sld xmlns:p=\"{P[1:-1]}\">{_timing(
        _STEP_TEMPLATE.format(behaviors=_behavior('5')),
        _STEP_TEMPLATE.format(behaviors=_behavior('6')),
    )}</p:sld>"
    steps = parse_steps(ET.fromstring(xml))
    assert steps == [[Target("5")], [Target("6")]]


def test_parse_steps_dedupes_repeated_behaviors(tmp_path):
    """同一个目标常被多条 cBhvr 重复引用（淡入 + 位移是两条），
    真实课件里 slide29 一步就有 2 条 cBhvr 指同一个 spid。"""
    xml = f"<p:sld xmlns:p=\"{P[1:-1]}\">{_timing(
        _STEP_TEMPLATE.format(behaviors=_behavior('6') + _behavior('6')),
    )}</p:sld>"
    assert parse_steps(ET.fromstring(xml)) == [[Target("6")]]


def test_parse_steps_reads_paragraph_range(tmp_path):
    xml = f"<p:sld xmlns:p=\"{P[1:-1]}\">{_timing(
        _STEP_TEMPLATE.format(behaviors=_behavior('3', para=0)),
        _STEP_TEMPLATE.format(behaviors=_behavior('3', para=1)),
    )}</p:sld>"
    steps = parse_steps(ET.fromstring(xml))
    assert steps == [[Target("3", 0, 0)], [Target("3", 1, 1)]]


def test_parse_steps_ignores_non_entrance(tmp_path):
    """强调/退出/路径动画不改变可见性，对分页没有意义。"""
    step = _STEP_TEMPLATE.format(behaviors=_behavior("5")).replace(
        'presetClass="entr"', 'presetClass="exit"'
    )
    xml = f'<p:sld xmlns:p="{P[1:-1]}">{_timing(step)}</p:sld>'
    assert parse_steps(ET.fromstring(xml)) == []


def test_parse_steps_empty_without_timing(tmp_path):
    assert parse_steps(ET.fromstring(f'<p:sld xmlns:p="{P[1:-1]}"/>')) == []


# ---- 展开 ----


def test_no_animation_leaves_file_untouched(tmp_path):
    deck = _deck(tmp_path / "a.pptx")
    before = deck.read_bytes()

    result = expand_animations(deck)

    assert result.expanded is False
    assert deck.read_bytes() == before


def test_one_click_becomes_two_pages(tmp_path):
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base,
        tmp_path / "a.pptx",
        1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "后出现的框"),
    )

    result = expand_animations(deck)

    assert result.expanded is True
    assert result.pages_before == 1
    assert result.pages_after == 2
    assert _pages(deck) == 2


def test_expanded_version_hides_the_animated_shape(tmp_path):
    """第一版必须看不见那个要点击才出现的形状——否则展开就没有意义。"""
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base,
        tmp_path / "a.pptx",
        1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "后出现的框"),
    )

    expand_animations(deck)

    # 新增的那一版排在原页之前
    assert "99" not in _slide_shape_ids(deck, "ppt/slides/slide2.xml")
    # 原页保留全部内容
    assert "99" in _slide_shape_ids(deck, "ppt/slides/slide1.xml")


def test_two_clicks_reveal_progressively(tmp_path):
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base,
        tmp_path / "a.pptx",
        1,
        _timing(
            _STEP_TEMPLATE.format(behaviors=_behavior("101")),
            _STEP_TEMPLATE.format(behaviors=_behavior("102")),
        ),
        shapes_xml=_shape("101", "一") + _shape("102", "二"),
    )

    result = expand_animations(deck)

    assert result.pages_after == 3
    v0 = _slide_shape_ids(deck, "ppt/slides/slide2.xml")   # 一个都没出现
    v1 = _slide_shape_ids(deck, "ppt/slides/slide3.xml")   # 出现了第一个
    v2 = _slide_shape_ids(deck, "ppt/slides/slide1.xml")   # 原页，全部出现

    assert "101" not in v0 and "102" not in v0
    assert "101" in v1 and "102" not in v1
    assert "101" in v2 and "102" in v2


def test_paragraph_animation_reveals_lines_one_by_one(tmp_path):
    """用户抱怨的「所有元素叠在一页上」，最常见的来源就是逐条出现的列表。"""
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base,
        tmp_path / "a.pptx",
        1,
        _timing(
            _STEP_TEMPLATE.format(behaviors=_behavior("77", para=0)),
            _STEP_TEMPLATE.format(behaviors=_behavior("77", para=1)),
        ),
        shapes_xml=_shape("77", "要点", paragraphs=3),
    )

    expand_animations(deck)

    def paras(part):
        with zipfile.ZipFile(deck) as zf:
            root = ET.fromstring(zf.read(part))
        for sp in root.iter(f"{P}sp"):
            c = sp.find(f".//{P}cNvPr")
            if c is not None and c.get("id") == "77":
                body = sp.find(f".//{P}txBody")
                return len(body.findall("{http://schemas.openxmlformats.org/drawingml/2006/main}p"))
        return None

    assert paras("ppt/slides/slide2.xml") == 1   # 只剩没有动画的第 3 段
    assert paras("ppt/slides/slide3.xml") == 2
    assert paras("ppt/slides/slide1.xml") == 3


def test_result_is_a_valid_pptx(tmp_path):
    """展开后的包必须仍是合法 pptx——五期的教训：改完 OPC 结构不验证
    「还能不能被打开」，等于什么都没做。"""
    base = _deck(tmp_path / "base.pptx", pages=2)
    deck = _inject(
        base,
        tmp_path / "a.pptx",
        1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "框"),
    )

    expand_animations(deck)

    prs = Presentation(str(deck))
    assert len(prs.slides) == 3


def test_no_dangling_relationships(tmp_path):
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base, tmp_path / "a.pptx", 1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "框"),
    )

    expand_animations(deck)

    from app.services.opc_rewrite import owner_part, read_rels

    with zipfile.ZipFile(deck) as zf:
        names = set(zf.namelist())
        for name in list(names):
            if not name.endswith(".rels"):
                continue
            for _rid, _type, target in read_rels(zf, owner_part(name)):
                assert target in names, f"{name} 指向不存在的 {target}"


def test_content_types_covers_new_slides(tmp_path):
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base, tmp_path / "a.pptx", 1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "框"),
    )

    expand_animations(deck)

    with zipfile.ZipFile(deck) as zf:
        ct = zf.read("[Content_Types].xml").decode("utf-8")
    assert "/ppt/slides/slide2.xml" in ct


# ---- 上限与跳过 ----


def test_too_many_steps_is_skipped_with_warning(tmp_path):
    """做了几十步动画的一页会把 deck 变成上千页，转换必然超时，
    而对做笔记毫无用处。"""
    base = _deck(tmp_path / "base.pptx", pages=1)
    steps = [
        _STEP_TEMPLATE.format(behaviors=_behavior(str(200 + i)))
        for i in range(MAX_STEPS_PER_SLIDE + 1)
    ]
    shapes = "".join(_shape(str(200 + i), f"s{i}") for i in range(MAX_STEPS_PER_SLIDE + 1))
    deck = _inject(base, tmp_path / "a.pptx", 1, _timing(*steps), shapes_xml=shapes)

    result = expand_animations(deck)

    assert result.expanded is False
    assert any("超过上限" in w for w in result.warnings)


def test_alternate_content_page_is_skipped_with_warning(tmp_path):
    """AlternateContent 里有 mc:Choice/mc:Fallback 两套并存的内容，
    改写它等于赌哪一套生效。跳过并告知，不假装成功。"""
    base = _deck(tmp_path / "base.pptx", pages=1)
    shapes = '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006"><mc:Choice Requires="a14"/></mc:AlternateContent>'
    deck = _inject(
        base, tmp_path / "a.pptx", 1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=shapes + _shape("99", "框"),
    )

    result = expand_animations(deck)

    assert result.expanded is False
    assert any("AlternateContent" in w for w in result.warnings)


def test_timing_is_dropped_from_every_version(tmp_path):
    """五期教训的重演：删了形状却留着引用它的 <p:timing>，Office 在线服务
    整份返回 406——而 LibreOffice 照转不误，所以只测 LibreOffice 会假绿。
    真机上这条正是第一次验证时炸出来的。
    """
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base, tmp_path / "a.pptx", 1,
        _timing(_STEP_TEMPLATE.format(behaviors=_behavior("99"))),
        shapes_xml=_shape("99", "框"),
    )

    expand_animations(deck)

    with zipfile.ZipFile(deck) as zf:
        for name in zf.namelist():
            if name.startswith("ppt/slides/slide"):
                assert b"<p:timing>" not in zf.read(name), f"{name} 仍留着 timing"


def test_no_animation_target_dangles(tmp_path):
    """更强的表述：任何一版里都不该存在「动画指向不存在的形状」。"""
    base = _deck(tmp_path / "base.pptx", pages=1)
    deck = _inject(
        base, tmp_path / "a.pptx", 1,
        _timing(
            _STEP_TEMPLATE.format(behaviors=_behavior("101")),
            _STEP_TEMPLATE.format(behaviors=_behavior("102")),
        ),
        shapes_xml=_shape("101", "一") + _shape("102", "二"),
    )

    expand_animations(deck)

    with zipfile.ZipFile(deck) as zf:
        for name in zf.namelist():
            if not name.startswith("ppt/slides/slide"):
                continue
            root = ET.fromstring(zf.read(name))
            present = set(_slide_shape_ids(deck, name))
            for sp_tgt in root.iter(f"{P}spTgt"):
                assert sp_tgt.get("spid") in present, f"{name} 的动画指向已删除的形状"
