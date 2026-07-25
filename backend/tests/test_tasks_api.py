import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.util import Emu
from pypdf import PdfReader

from app.config import settings


@pytest.fixture
def client(tmp_path, monkeypatch):
    # Base/engine 延迟到函数体内导入，理由见 test_uploads_api.py 里同名 fixture
    # 的注释：conftest.py 的 _isolate_app_db autouse fixture 重定向了
    # app.db.engine，模块顶层 import 会绕过这个重定向。
    from app.db import Base, engine

    monkeypatch.setattr(settings, "storage_root", tmp_path / "storage")
    monkeypatch.setattr(settings, "chunk_size", 64 * 1024)
    settings.ensure_dirs()
    Base.metadata.drop_all(engine)
    Base.metadata.create_all(engine)

    from app.main import app

    return AsyncClient(transport=ASGITransport(app=app), base_url="http://test")


@pytest.fixture
def sample_bytes(tmp_path):
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "s.pptx"
    prs.save(path)
    return path.read_bytes()


async def _upload(client, payload: bytes) -> str:
    created = (
        await client.post(
            "/api/uploads", json={"filename": "deck.pptx", "size": len(payload)}
        )
    ).json()
    uid, size = created["upload_id"], created["chunk_size"]
    for idx in range(created["total_chunks"]):
        await client.put(
            f"/api/uploads/{uid}/chunks/{idx}",
            content=payload[idx * size : (idx + 1) * size],
        )
    return (await client.post(f"/api/uploads/{uid}/complete")).json()["task_id"]


async def test_task_reaches_done_with_parsed_meta(client, sample_bytes):
    task_id = await _upload(client, sample_bytes)

    body = (await client.get(f"/api/tasks/{task_id}")).json()
    assert body["status"] == "done"
    assert body["slide_count"] == 4
    assert body["slide_width_emu"] == 12192000
    assert body["engine"] == "placeholder"
    assert isinstance(body["fonts"], list)


async def test_download_returns_pdf_with_matching_pages(client, sample_bytes, tmp_path):
    task_id = await _upload(client, sample_bytes)

    resp = await client.get(f"/api/tasks/{task_id}/download")
    assert resp.status_code == 200
    assert resp.headers["content-type"] == "application/pdf"

    out = tmp_path / "got.pdf"
    out.write_bytes(resp.content)
    assert len(PdfReader(out).pages) == 4


async def test_task_not_found(client):
    resp = await client.get("/api/tasks/nope")
    assert resp.status_code == 404


async def test_invalid_pptx_marks_failed(client):
    task_id = await _upload(client, b"definitely not a zip file")

    body = (await client.get(f"/api/tasks/{task_id}")).json()
    assert body["status"] == "failed"
    assert body["error_code"] == "PPTX_INVALID_ZIP"


async def test_download_before_done_returns_409(client):
    task_id = await _upload(client, b"definitely not a zip file")

    resp = await client.get(f"/api/tasks/{task_id}/download")
    assert resp.status_code == 409


async def test_run_task_walks_full_state_machine(client, sample_bytes):
    """回归测试：Global Constraint 规定任务状态机固定为
    pending -> parsing -> queued -> converting -> done，占位引擎瞬时完成
    也必须走完全部状态——这条约束存在的唯一理由是二期接真引擎时前端不用改。

    但 grep 全部测试，"queued" 和 "converting" 从未作为断言目标出现过：
    把 pipeline.py 里那两行 _set_status(session, task, "queued") 和
    "converting" 直接删掉，其余 44 个测试依然全绿。这里用 SQLAlchemy 的
    attribute event 直接观测 Task.status 的赋值序列，堵住这个洞。
    """
    from sqlalchemy import event

    from app.db import SessionLocal
    from app.models import Task
    from app.services import pipeline

    task_id = "task-state-machine-sequence"
    src = settings.originals_dir / f"{task_id}.pptx"
    src.write_bytes(sample_bytes)

    setup = SessionLocal()
    setup.add(
        Task(
            task_id=task_id,
            upload_id="upload-does-not-matter",
            original_filename="deck.pptx",
            size_bytes=len(sample_bytes),
            status="pending",
            engine="placeholder",
        )
    )
    setup.commit()
    setup.close()

    transitions: list[str] = []

    def _record(target, value, oldvalue, initiator):
        if getattr(target, "task_id", None) == task_id:
            transitions.append(value)

    event.listen(Task.status, "set", _record)
    try:
        pipeline.run_task(task_id)
    finally:
        event.remove(Task.status, "set", _record)

    assert transitions == ["parsing", "queued", "converting", "done"]
