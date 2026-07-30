import zipfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services import media_strip
from app.services.media_strip import MEDIA_REL_TYPES, strip_media
from app.services.opc_rewrite import read_rels

VIDEO_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
MEDIA_REL = "http://schemas.microsoft.com/office/2007/relationships/media"


def _plain_deck(path: Path, pages: int = 3) -> Path:
    """不含媒体的普通 deck。"""
    prs = Presentation()
    for i in range(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"PAGE-{i + 1}"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    prs.save(str(path))
    return path


IMAGE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"

# 1x1 透明 PNG，当海报帧用
_TINY_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000a49444154789c630001000005000105"
    "0d0a2db40000000049454e44ae426082"
)

# 真实 PowerPoint 嵌视频后 slide 正文里的结构，逐字照抄自用户那份 83.7MB
# 课件的 slide25（rId 换成本 fixture 用的名字）。四处引用缺一不可——只删
# part 和 Relationship 而留着这些，Office 在线服务会整份拒绝（HTTP 406）。
_PIC_WITH_VIDEO = (
    '<p:pic><p:nvPicPr>'
    '<p:cNvPr id="5" name="Video">'
    '<a:hlinkClick r:id="" action="ppaction://media"/>'
    "</p:cNvPr>"
    '<p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr>'
    "<p:nvPr>"
    '<a:videoFile r:link="rIdVid"/>'
    "<p:extLst>"
    '<p:ext uri="{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}">'
    '<p14:media xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"'
    ' r:embed="rIdMed"/>'
    "</p:ext>"
    "</p:extLst>"
    "</p:nvPr>"
    "</p:nvPicPr>"
    '<p:blipFill><a:blip r:embed="rIdPoster"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
    '<p:spPr><a:xfrm><a:off x="954113" y="2089033"/><a:ext cx="6353781" cy="3574000"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
    "</p:pic>"
)

# <p:timing> 里的 <p:video> 用 spid 指形状、不带 r:id，所以删关系不会让它
# 「悬空」——但它声明「这个形状是媒体」，媒体没了就前后矛盾。真机实测：
# 不删这一段，即使前三处都清干净了，Office 仍然 406。
_TIMING_WITH_VIDEO = (
    "<p:timing><p:tnLst><p:par><p:cTn id=\"1\" dur=\"indefinite\" restart=\"never\""
    ' nodeType="tmRoot"><p:childTnLst>'
    '<p:video><p:cMediaNode vol="80000"><p:cTn id="7" fill="hold" display="0">'
    '<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst></p:cTn>'
    '<p:tgtEl><p:spTgt spid="5"/></p:tgtEl></p:cMediaNode></p:video>'
    "</p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"
)


def _deck_with_fake_video(path: Path, payload_mb: int = 2) -> Path:
    """造一个带假视频的 deck，**结构与真实 PowerPoint 产出一致**。

    python-pptx 没有加视频的稳定 API，所以先造普通 deck，再手工补齐四样：
    media part、slide1 rels 上的 video + media 两条关系（PowerPoint 真实
    产出就是两条都写）、`[Content_Types].xml` 的 Default，**以及 slide1
    正文里引用它们的 `<p:pic>` 与 `<p:timing>`**。

    最后那一项是这个 fixture 的要害：本文件早先的版本只注入了关系而正文
    里没有任何引用，剥离后自然不会留下死引用——测的是真实世界不存在的
    形状，于是「剥离后能转」全是假绿，真实课件一上传就 406。
    """
    tmp = path.with_suffix(".base.pptx")
    _plain_deck(tmp, pages=3)

    payload = b"\x00" * (payload_mb * 1024 * 1024)
    with zipfile.ZipFile(tmp) as zin, zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            raw = zin.read(item.filename)
            if item.filename == "ppt/slides/_rels/slide1.xml.rels":
                injected = (
                    f'<Relationship Id="rIdVid" Type="{VIDEO_REL}" Target="../media/movie1.mp4"/>'
                    f'<Relationship Id="rIdMed" Type="{MEDIA_REL}" Target="../media/movie1.mp4"/>'
                    f'<Relationship Id="rIdPoster" Type="{IMAGE_REL}" Target="../media/poster1.png"/>'
                ).encode()
                raw = raw.replace(b"</Relationships>", injected + b"</Relationships>")
            elif item.filename == "ppt/slides/slide1.xml":
                text = raw.decode("utf-8")
                text = text.replace("</p:spTree>", _PIC_WITH_VIDEO + "</p:spTree>")
                text = text.replace("</p:sld>", _TIMING_WITH_VIDEO + "</p:sld>")
                raw = text.encode("utf-8")
            elif item.filename == "[Content_Types].xml":
                raw = raw.replace(
                    b"</Types>",
                    b'<Default Extension="mp4" ContentType="video/mp4"/>'
                    b'<Default Extension="png" ContentType="image/png"/></Types>',
                )
            zout.writestr(item, raw)
        zout.writestr("ppt/media/movie1.mp4", payload)
        zout.writestr("ppt/media/poster1.png", _TINY_PNG)

    tmp.unlink()
    return path


def _parts(path: Path) -> set[str]:
    with zipfile.ZipFile(path) as zf:
        return set(zf.namelist())


def test_media_rel_types_covers_video_audio_and_media():
    # PowerPoint 嵌一段视频会同时写 video 与 media 两条关系指向同一个 part，
    # 漏一个文件就还留在包里
    assert VIDEO_REL in MEDIA_REL_TYPES
    assert MEDIA_REL in MEDIA_REL_TYPES
    assert any("audio" in t for t in MEDIA_REL_TYPES)


def test_strips_video_part(tmp_path):
    deck = _deck_with_fake_video(tmp_path / "v.pptx", payload_mb=2)
    before = _parts(deck)
    assert "ppt/media/movie1.mp4" in before

    result = strip_media(deck)

    assert result.stripped is True
    assert result.removed_parts == 1
    assert result.bytes_after < result.bytes_before
    assert "ppt/media/movie1.mp4" not in _parts(deck)


def _slide1(path: Path) -> str:
    with zipfile.ZipFile(path) as zf:
        return zf.read("ppt/slides/slide1.xml").decode("utf-8")


def test_fixture_really_references_the_media_in_slide_body(tmp_path):
    """守住 fixture 本身：正文里必须真的引用了媒体关系。

    这个 fixture 早先只往 .rels 注入关系、正文里没有任何引用，于是剥离后
    不可能留下死引用——「剥离后能转」因此全是假绿，直到真实课件上传后
    Office 服务返回 406 才暴露。这条测试钉住 fixture 的保真度：它一旦退化
    回那个形状，下面几条断言就都失去意义。
    """
    xml = _slide1(_deck_with_fake_video(tmp_path / "v.pptx"))
    assert 'r:link="rIdVid"' in xml
    assert 'r:embed="rIdMed"' in xml
    assert 'action="ppaction://media"' in xml
    assert "<p:video>" in xml


def test_strips_video_reference_from_slide_body(tmp_path):
    """真机根因：光删 part 和 Relationship，Office 服务会整份拒绝（406）。

    四处引用逐条断言——真机对照实验证明它们缺一不可：只清前三处而留下
    <p:timing> 里的 <p:video>，Graph 仍然 406。
    """
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    xml = _slide1(deck)

    assert "<a:videoFile" not in xml
    assert "p14:media" not in xml
    assert "ppaction://media" not in xml
    assert "<p:video>" not in xml


def test_keeps_poster_frame_so_the_slide_is_not_blank(tmp_path):
    """<p:pic> 与它的海报帧一律保留：摘掉媒体引用后它退化成一张普通图片，
    那一页显示视频封面而不是空白，信息损失最小。"""
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    xml = _slide1(deck)

    assert "<p:pic>" in xml
    assert 'r:embed="rIdPoster"' in xml
    assert "ppt/media/poster1.png" in _parts(deck)


def test_body_cleanup_leaves_no_empty_shells(tmp_path):
    """摘掉 p14:media 那个 <p:ext> 后不留空的 <p:extLst> 残壳。"""
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    assert "<p:extLst></p:extLst>" not in _slide1(deck)


def test_slides_without_media_are_copied_verbatim(tmp_path):
    """只有挂着媒体关系的 part 才做正文改写，其余逐字节复制。

    改写面越小，撞上 mc:Ignorable 被 ET 往返吃掉那类坑的机会越小——
    三期在这上面修了五轮。
    """
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    with zipfile.ZipFile(deck) as zf:
        before = {n: zf.read(n) for n in zf.namelist() if n.startswith("ppt/slides/slide")}
    strip_media(deck)
    with zipfile.ZipFile(deck) as zf:
        for name, raw in before.items():
            if name.endswith("slide1.xml"):
                continue  # 唯一带媒体的那页，本来就该被改写
            assert zf.read(name) == raw, f"{name} 不该被改动"


def test_keeps_all_slides(tmp_path):
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    slides = [n for n in _parts(deck) if n.startswith("ppt/slides/slide")]
    assert len(slides) == 3


def test_no_dangling_internal_relationships(tmp_path):
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    with zipfile.ZipFile(deck) as zf:
        names = set(zf.namelist())
        for name in list(names):
            if not name.endswith(".rels"):
                continue
            owner = name.replace("/_rels/", "/").replace(".rels", "")
            for _rid, _type, target in read_rels(zf, owner):
                assert target in names, f"{name} 指向不存在的 {target}"


def test_result_is_openable(tmp_path):
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    strip_media(deck)
    prs = Presentation(str(deck))
    assert len(prs.slides) == 3


def test_deck_without_media_is_not_rewritten(tmp_path):
    deck = _plain_deck(tmp_path / "p.pptx")
    before_bytes = deck.read_bytes()

    result = strip_media(deck)

    assert result.stripped is False
    assert result.removed_parts == 0
    # 零收益时不该白做一次解压重打包
    assert deck.read_bytes() == before_bytes


def test_preserves_mc_ignorable(tmp_path):
    """ET 往返会丢弃未使用的 xmlns 声明，而 mc:Ignorable 指着它们。
    presentation.xml 必须逐字节保留——这是三期修了一整轮的坑。"""
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    with zipfile.ZipFile(deck) as zf:
        pres_before = zf.read("ppt/presentation.xml")

    strip_media(deck)

    with zipfile.ZipFile(deck) as zf:
        assert zf.read("ppt/presentation.xml") == pres_before


def test_write_failure_cleans_up_tmp_and_preserves_source(tmp_path, monkeypatch):
    """写入过程中失败时，临时文件要被清理，源文件不能被破坏。

    originals_dir 目前没有其它机制会兜底清理孤儿 tmp 文件——drop_original
    只精确删 {task_id}.pptx，purge_expired_outputs/purge_expired_shards 都不
    扫这里——所以 strip_media 的 except 分支必须自己证明语义正确：异常穿透
    给调用方，tmp 文件不留下，源文件全程未被替换。
    """
    deck = _deck_with_fake_video(tmp_path / "v.pptx")
    before_bytes = deck.read_bytes()

    def _boom(*_args, **_kwargs):
        raise RuntimeError("boom")

    monkeypatch.setattr(media_strip, "rewrite_content_types", _boom)

    with pytest.raises(RuntimeError, match="boom"):
        strip_media(deck)

    # 源文件全程未被 tmp.replace(src) 替换，字节不变
    assert deck.read_bytes() == before_bytes
    # 没有留下孤儿 tmp 文件
    leftover = [p for p in tmp_path.iterdir() if p != deck and p.suffix == ".pptx"]
    assert leftover == []
