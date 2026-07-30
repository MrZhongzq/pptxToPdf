import pytest
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Emu

from app.config import settings


@pytest.fixture
def client(tmp_path, monkeypatch):
    # Base/engine 延迟到函数体内导入，理由见 test_uploads_api.py 里同名 fixture
    # 的注释：conftest.py 的 _isolate_app_db autouse fixture 重定向了
    # app.db.engine，模块顶层 import 会绕过这个重定向。
    from app.db import Base, engine

    monkeypatch.setattr(settings, "storage_root", tmp_path / "storage")
    monkeypatch.setattr(settings, "chunk_size", 64 * 1024)
    settings.ensure_dirs()
    Base.metadata.drop_all(engine)
    Base.metadata.create_all(engine)

    from app.main import app

    with TestClient(app) as c:
        yield c


def _sample_deck_bytes() -> bytes:
    import io

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for _ in range(4):
        prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def logged_in(client, monkeypatch):
    """六期起 Graph 通道要求登录。用 Graph 的用例都要先过这道门。"""
    from cryptography.fernet import Fernet

    import app.db as db_module
    from app.services import auth, users

    monkeypatch.setattr(auth.settings, "secret_key", Fernet.generate_key().decode())
    monkeypatch.setattr(auth.settings, "admin_cookie_secure", False)
    monkeypatch.setattr(auth, "_WRONG_PASSWORD_DELAY_S", 0.0)
    with db_module.SessionLocal() as db:
        users.create(db, username="alice", email="a@e.com", password="hunter2!")
    resp = client.post("/api/auth/login", json={"username": "alice", "password": "hunter2!"})
    assert resp.status_code == 200, resp.text
    return client


@pytest.fixture
def db_session():
    import app.db as db_module

    s = db_module.SessionLocal()
    yield s
    s.close()


def _upload_a_deck(client, engine: str | None = None) -> str:
    """走完整的分块上传协议，落地一个 ready 状态的任务，返回 task_id。

    engine：透传给 POST /api/uploads 的 `engine` 字段（用户上传时选的引擎，
    未废弃，参见 CreateUploadRequest）。默认 None，即不指定。
    """
    payload = _sample_deck_bytes()
    body = {"filename": "deck.pptx", "size": len(payload)}
    if engine is not None:
        body["engine"] = engine
    created = client.post("/api/uploads", json=body).json()
    uid, size = created["upload_id"], created["chunk_size"]
    for idx in range(created["total_chunks"]):
        client.put(
            f"/api/uploads/{uid}/chunks/{idx}",
            content=payload[idx * size : (idx + 1) * size],
        )
    return client.post(f"/api/uploads/{uid}/complete").json()["task_id"]


def _session():
    import app.db as db_module

    return db_module.SessionLocal()


def _get_task(task_id: str) -> dict:
    """走真实的 get_task 端点函数（不经 HTTP 层），拿到与 API 响应同构的 dict。"""
    from app.api.tasks import get_task as _get_task_endpoint

    session = _session()
    try:
        return _get_task_endpoint(task_id, session).model_dump()
    finally:
        session.close()


def _load_task_row(task_id: str):
    from app.models import Task

    session = _session()
    try:
        return session.get(Task, task_id)
    finally:
        session.close()


def test_complete_leaves_task_ready_and_does_not_enqueue(client, monkeypatch):
    """complete 只拼装落库，不入队。

    fix round：uploads.py 已经不再 import enqueue_conversion（谁在
    complete_upload 里写 enqueue_conversion(task_id) 现在会直接 NameError，
    比一条测试更硬）。monkeypatch.setattr 默认要求属性已存在，这里用
    raising=False 让它在 app.api.uploads 的模块 globals 上新建这个名字——
    如果生产代码真被改回去调用它，运行时依然会解析到这个假替身并被
    enqueued 断言抓住；不加 raising=False 则会在 setattr 这一步就
    AttributeError，反而测不到「加回调用」这个变异。
    """
    enqueued = []
    monkeypatch.setattr(
        "app.api.uploads.enqueue_conversion",
        lambda t: enqueued.append(t),
        raising=False,
    )
    task_id = _upload_a_deck(client)
    task = _get_task(task_id)
    assert task["status"] == "ready"
    assert enqueued == [], "complete 不该入队"


def test_complete_upload_lazily_purges_expired_ready_tasks(client):
    """终审 M-1：purge_expired_ready 原来的两个触发点是「API 启动」和
    「任意转换跑完」，但 ready 状态的原文件是"上传"产生的，不是"转换"
    产生的。这里有个自己够不到自己的闭环：ready TTL 存在的全部目的就是
    防"用户只传不点开始"，而在那个场景下恰恰没有任何转换会跑完——回收器
    因此永不触发，README 说的"最坏晚不少"实际是"最坏永不"。

    complete_upload 才是"上传"这个动作真正发生的地方，必须在这里顺带
    触发一次。造一个已经过期的 ready 任务 T-old（不经过 /start，模拟
    "只传不点开始"），再完整走一次上传流程触发一次新的 complete_upload
    ——T-old 必须被顺带回收，不必等任何转换跑完或服务重启。
    """
    from datetime import datetime, timedelta, timezone

    from app.models import Task

    old_task_id = _upload_a_deck(client)  # 落 ready，不 start——正是回收器
    # 原本永远够不到自己的那个场景。

    session = _session()
    try:
        row = session.get(Task, old_task_id)
        row.updated_at = datetime.now(timezone.utc).replace(tzinfo=None) - timedelta(
            hours=settings.ready_ttl_hours, minutes=1
        )
        session.commit()
    finally:
        session.close()

    _upload_a_deck(client)  # 触发新一次 complete_upload，顺带回收 T-old

    old_row = _load_task_row(old_task_id)
    assert old_row.status == "failed"
    assert old_row.error_code == "READY_EXPIRED"


def test_start_enqueues_and_moves_to_pending(client, monkeypatch):
    enqueued = []
    monkeypatch.setattr("app.api.tasks.enqueue_conversion", lambda t: enqueued.append(t))
    task_id = _upload_a_deck(client)
    resp = client.post(f"/api/tasks/{task_id}/start", json={"engine": "libreoffice"})
    assert resp.status_code == 200
    assert resp.json()["status"] == "pending"
    assert enqueued == [task_id]


def test_start_records_engine_and_options(client, logged_in, monkeypatch):
    """引擎与选项在 start 时才定——这正是本期的目的。"""
    monkeypatch.setattr("app.api.tasks.enqueue_conversion", lambda t: None)
    task_id = _upload_a_deck(client)
    resp = client.post(
        f"/api/tasks/{task_id}/start",
        json={"engine": "graph", "options": {"expand_animations": True}},
    )
    assert resp.status_code == 200, resp.text
    task = _load_task_row(task_id)
    assert task.requested_engine == "graph"
    assert "expand_animations" in (task.options_json or "")


def test_start_without_engine_keeps_the_one_chosen_at_upload(client, logged_in, monkeypatch):
    """fix round I1：上传时选了引擎、start 不带 engine 时不该被静默清空。

    complete_upload 把 upload.requested_engine 转写进 task.requested_engine；
    start 之前它是无条件覆盖成 payload.engine（默认 None），上传时选的引擎
    会在没人碰它的情况下消失。用户裁决沿用上传时选的——start 只在
    payload.engine 非 None 时才覆盖。
    """
    monkeypatch.setattr("app.api.tasks.enqueue_conversion", lambda t: None)
    task_id = _upload_a_deck(client, engine="graph")
    resp = client.post(f"/api/tasks/{task_id}/start", json={})
    # 断言状态码，不只看 DB：六期加 Graph 登录门槛时发现这条测试原本
    # 在 401 下也能"通过"——请求被拦、UPDATE 没执行，而 requested_engine
    # 还是 complete 时写的 graph，断言碰巧成立。少一条状态码断言，一条
    # 鉴权回归就能从它眼皮底下溜过去。
    assert resp.status_code == 200, resp.text
    task = _load_task_row(task_id)
    assert task.requested_engine == "graph"


def test_start_with_explicit_engine_overrides_the_one_chosen_at_upload(client, logged_in, monkeypatch):
    """终审 I-3：现有两条测试各只覆盖了一半——「上传没选/start 选了」
    （test_start_records_engine_and_options）和「上传选了/start 没选」
    （test_start_without_engine_keeps_the_one_chosen_at_upload）——恰好都
    躲开了「上传选 A、start 显式选 B」这条路径。终审做过变异实验：把
    `if payload.engine is not None:` 误改成 `if task.requested_engine is
    None:`（典型的"条件填充"误实现——只在字段为空时才填充，而不是"传了就
    覆盖"），这两条既有测试依然全绿，因为它们都没有在"两边都有值"的情况下
    断言"新值赢"。这个误实现的症状正是项目第一铁律要挡的那种：用户在
    就绪卡片上把引擎从 libreoffice 改成 graph，后端却照旧用上传时选的
    libreoffice——两段式上传存在的全部理由就是让用户能在这一步改主意。
    """
    monkeypatch.setattr("app.api.tasks.enqueue_conversion", lambda t: None)
    task_id = _upload_a_deck(client, engine="libreoffice")
    client.post(f"/api/tasks/{task_id}/start", json={"engine": "graph"})
    task = _load_task_row(task_id)
    assert task.requested_engine == "graph"


def test_start_twice_is_409(client, monkeypatch):
    enqueued = []
    monkeypatch.setattr("app.api.tasks.enqueue_conversion", lambda t: enqueued.append(t))
    task_id = _upload_a_deck(client)
    client.post(f"/api/tasks/{task_id}/start", json={})
    resp = client.post(f"/api/tasks/{task_id}/start", json={})
    assert resp.status_code == 409
    assert resp.json()["code"] == "TASK_ALREADY_STARTED"
    assert len(enqueued) == 1, "不该重复入队"


def test_start_on_already_reaped_ready_task_is_410(client, monkeypatch):
    """purge_expired_ready 已经把这个任务标 failed + READY_EXPIRED 之后，
    /start 必须把「已被回收」和「真的已经在跑」（test_start_twice_is_409）
    区分开——前者是 410，不是笼统的 409。message 复用回收器自己写的那句，
    不在 start_task 里另起一份，两处措辞才不会跑偏。"""
    from datetime import datetime, timedelta, timezone

    from app.models import Task
    from app.services.retention import purge_expired_ready

    task_id = _upload_a_deck(client)

    session = _session()
    try:
        row = session.get(Task, task_id)
        row.updated_at = datetime.now(timezone.utc).replace(tzinfo=None) - timedelta(
            hours=settings.ready_ttl_hours, minutes=1
        )
        session.commit()
    finally:
        session.close()

    assert purge_expired_ready() == 1
    expected_message = _load_task_row(task_id).error_message

    resp = client.post(f"/api/tasks/{task_id}/start", json={})
    assert resp.status_code == 410
    assert resp.json()["code"] == "READY_EXPIRED"
    assert resp.json()["message"] == expected_message


def test_start_on_missing_task_is_404(client):
    resp = client.post("/api/tasks/does-not-exist/start", json={})
    assert resp.status_code == 404
    assert resp.json()["code"] == "TASK_NOT_FOUND"


def test_concurrent_start_only_wins_once(client, monkeypatch):
    """终审 I-4：start_task 原本是"读判改写"三步走——读 status、判断、
    单独一次 commit 写 "pending"，两步之间有窗口。FastAPI 把同步端点丢进
    线程池，两个并发 /start 请求都可能在窗口内读到 status=='ready'，各自
    以为自己有权推进。终审用 TestClient 两线程同打同一个 ready 任务、连打
    三轮，第二轮两个请求都拿到 200、同一个 task_id 被入队两次——这不只是
    浪费资源，两个 run_task 并发跑同一个 task_id 会同时改写同一份原文件、
    同时写同一个输出 PDF，终态还会在 done/failed 之间乱翻转。

    只用 threading.Barrier 卡两个请求的发出时刻不够可靠：本机 SQLite
    单条语句在微秒级完成，GIL 未必真的会在两个线程之间切换，实测这样写
    的测试即使打在旧实现上也会"侥幸通过"（本机反复验证：均为一读一写、
    从未真正交错）。这里改为在 `_load`（两种实现都共享的读入口）里插一个
    屏障，强制两个线程都读到 task.status=='ready' 之后才放行——这正是
    终审描述的那条窗口，不依赖 OS 线程调度的运气，实测能稳定把旧实现
    打成双 200（复现变异见 final-fix-report）。新实现下不管调度怎么交错，
    条件 UPDATE（WHERE status='ready'）保证只有一个请求能把 rowcount 从
    0 改到 1，结果必须是恰好一个 200、一个 409，且入队恰好一次。
    """
    import threading

    import app.api.tasks as tasks_module

    original_load = tasks_module._load
    read_barrier = threading.Barrier(2, timeout=5)

    def synced_load(session, task_id):
        task = original_load(session, task_id)
        read_barrier.wait()
        return task

    monkeypatch.setattr(tasks_module, "_load", synced_load)

    enqueue_lock = threading.Lock()
    enqueued: list[str] = []

    def _enqueue(t: str) -> None:
        with enqueue_lock:
            enqueued.append(t)

    monkeypatch.setattr("app.api.tasks.enqueue_conversion", _enqueue)
    task_id = _upload_a_deck(client)

    results_lock = threading.Lock()
    results: list[int] = []

    def fire():
        resp = client.post(f"/api/tasks/{task_id}/start", json={})
        with results_lock:
            results.append(resp.status_code)

    t1 = threading.Thread(target=fire)
    t2 = threading.Thread(target=fire)
    t1.start()
    t2.start()
    t1.join()
    t2.join()

    assert sorted(results) == [200, 409]
    assert enqueued == [task_id], "只能入队一次"
    assert _load_task_row(task_id).status == "pending"


def test_start_drops_original_when_enqueue_fails(client, monkeypatch):
    """Redis 挂了时必须删原文件——这段兜底是从 complete 挪过来的，
    漏挪的话每次 Redis 抖动都留一份 80-500MB 的孤儿。"""
    dropped = []
    monkeypatch.setattr("app.api.tasks.drop_original", lambda t: dropped.append(t))

    def boom(_):
        raise ConnectionError("redis down")

    monkeypatch.setattr("app.api.tasks.enqueue_conversion", boom)

    task_id = _upload_a_deck(client)
    resp = client.post(f"/api/tasks/{task_id}/start", json={})
    assert resp.status_code == 503
    assert dropped == [task_id]
    assert _load_task_row(task_id).status == "failed"


# ---- 六期：Graph 通道的登录门槛 ----


def _login_as(client, db, username="alice", role="user"):
    from app.services import users

    users.create(db, username=username, email=f"{username}@e.com", password="hunter2!", role=role)
    resp = client.post("/api/auth/login", json={"username": username, "password": "hunter2!"})
    assert resp.status_code == 200, resp.text
    return client


def test_anonymous_cannot_start_graph_task(client, monkeypatch):
    """前端把 Graph 选项置灰只是体验，这里才是边界——绕过前端直接打 API
    是最基本的渗透手法。"""
    from cryptography.fernet import Fernet

    from app.services import auth

    monkeypatch.setattr(auth.settings, "secret_key", Fernet.generate_key().decode())
    task_id = _upload_a_deck(client, engine="graph")

    resp = client.post(f"/api/tasks/{task_id}/start", json={})

    assert resp.status_code == 401
    assert resp.json()["code"] == "AUTH_REQUIRED"


def test_anonymous_cannot_start_graph_via_payload_either(client, monkeypatch):
    """上传时选 libreoffice、start 时改成 graph——同样要拦住。
    只看 task.requested_engine 会漏掉这条路径。"""
    from cryptography.fernet import Fernet

    from app.services import auth

    monkeypatch.setattr(auth.settings, "secret_key", Fernet.generate_key().decode())
    task_id = _upload_a_deck(client)

    resp = client.post(f"/api/tasks/{task_id}/start", json={"engine": "graph"})

    assert resp.status_code == 401
    assert resp.json()["code"] == "AUTH_REQUIRED"


def test_anonymous_can_still_start_libreoffice_task(client, monkeypatch):
    """需求只说 Graph 对未登录用户灰色，没说未登录不能用站点。
    匿名访客仍可用 LibreOffice——这与「当前不开放注册」的设定一致。"""
    from cryptography.fernet import Fernet

    from app.services import auth

    monkeypatch.setattr(auth.settings, "secret_key", Fernet.generate_key().decode())
    task_id = _upload_a_deck(client, engine="libreoffice")

    resp = client.post(f"/api/tasks/{task_id}/start", json={})

    assert resp.status_code == 200


def test_logged_in_user_can_start_graph_task(client, db_session, monkeypatch):
    from cryptography.fernet import Fernet

    from app.services import auth

    monkeypatch.setattr(auth.settings, "secret_key", Fernet.generate_key().decode())
    monkeypatch.setattr(auth, "_WRONG_PASSWORD_DELAY_S", 0.0)
    _login_as(client, db_session)
    task_id = _upload_a_deck(client, engine="graph")

    resp = client.post(f"/api/tasks/{task_id}/start", json={})

    assert resp.status_code == 200
