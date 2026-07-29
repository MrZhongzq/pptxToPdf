import io
import posixpath
import random
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches

from app.services.pptx_probe import probe
from app.services.pptx_split import split_pptx

SLIDES = 8

# 与 pptx_probe.SLIDE_RE 同口径，测试里独立定义避免依赖另一个模块的
# 内部实现细节。
_SLIDE_PART_RE = re.compile(r"^ppt/slides/slide\d+\.xml$")

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_SLIDE_REL_TYPE = f"{_R_NS}/slide"


def _noisy_png(w=700, h=450) -> io.BytesIO:
    """随机噪声图压缩率低，单张几百 KB——用来验证 media 确实被裁掉。"""
    img = Image.new("RGB", (w, h))
    img.putdata(
        [
            (random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
            for _ in range(w * h)
        ]
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


@pytest.fixture
def deck(tmp_path) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
        s.shapes.add_picture(_noisy_png(), Inches(1), Inches(2), width=Inches(4))
    path = tmp_path / "deck.pptx"
    prs.save(path)
    return path


def _media_count(path: Path) -> int:
    with zipfile.ZipFile(path) as z:
        return len([n for n in z.namelist() if n.startswith("ppt/media/")])


def test_split_page_counts(deck, tmp_path):
    out = split_pptx(deck, [(1, 3), (4, 8)], tmp_path / "shards")
    assert len(out) == 2
    assert probe(out[0]).slide_count == 3
    assert probe(out[1]).slide_count == 5


def test_split_drops_unreferenced_media(deck, tmp_path):
    assert _media_count(deck) == SLIDES
    out = split_pptx(deck, [(1, 3), (4, 8)], tmp_path / "shards")
    # 关键断言：只从 sldIdLst 删除而不 drop_rel 的实现会让所有 media
    # 留在包里，这里就会看到 8 而不是 3/5。
    assert _media_count(out[0]) == 3
    assert _media_count(out[1]) == 5


def test_split_reduces_file_size(deck, tmp_path):
    out = split_pptx(deck, [(1, 4), (5, 8)], tmp_path / "shards")
    total = sum(p.stat().st_size for p in out)
    # 两片各占一半 media，合计应明显小于原包的 1.5 倍
    # （共享的 theme/master 会在每片里各留一份，所以不是精确的 1.0 倍）
    assert total < deck.stat().st_size * 1.5


def test_split_preserves_page_order_and_content(deck, tmp_path):
    """切片后每页的标题必须还是原来那一页的标题，顺序不能乱。"""
    out = split_pptx(deck, [(1, 3), (4, 8)], tmp_path / "shards")

    first = Presentation(str(out[0]))
    assert [s.shapes.title.text for s in first.slides] == [
        "PAGE-1",
        "PAGE-2",
        "PAGE-3",
    ]
    second = Presentation(str(out[1]))
    assert [s.shapes.title.text for s in second.slides] == [
        "PAGE-4",
        "PAGE-5",
        "PAGE-6",
        "PAGE-7",
        "PAGE-8",
    ]


def test_split_output_is_valid_pptx(deck, tmp_path):
    """产出的包必须能被 python-pptx 打开、被 pptx_probe 解析。"""
    out = split_pptx(deck, [(2, 5)], tmp_path / "shards")
    meta = probe(out[0])
    assert meta.slide_count == 4
    assert meta.slide_width_emu == 12192000
    prs = Presentation(str(out[0]))
    assert len(prs.slides) == 4


def test_single_page_range(deck, tmp_path):
    out = split_pptx(deck, [(5, 5)], tmp_path / "shards")
    assert probe(out[0]).slide_count == 1
    assert Presentation(str(out[0])).slides[0].shapes.title.text == "PAGE-5"


def test_full_range_is_a_faithful_copy(deck, tmp_path):
    out = split_pptx(deck, [(1, SLIDES)], tmp_path / "shards")
    assert probe(out[0]).slide_count == SLIDES
    assert _media_count(out[0]) == SLIDES


def test_rejects_out_of_bounds_range(deck, tmp_path):
    with pytest.raises(ValueError):
        split_pptx(deck, [(1, SLIDES + 1)], tmp_path / "shards")
    with pytest.raises(ValueError):
        split_pptx(deck, [(0, 3)], tmp_path / "shards")


# ---------------------------------------------------------------------------
# 代码审查追加的三个测试，见 task-3-report.md「代码审查追加修复」一节：
# 1) 分片内不能留下悬空 Relationship（不只是 presentation.xml.rels）；
# 2) 内部跳转（hlinkClick，Type 与 sldIdLst 里的 slide 关系相同）不能把
#    range 外的整页拖回分片；
# 3) presentation.xml 改成正则文本手术后，原有的 mc:Ignorable 与它对应
#    的 xmlns 声明必须原样保留，不能被 ET 往返丢弃。
# ---------------------------------------------------------------------------


def _owner_part(rels_name: str) -> str:
    """反向映射 pptx_split._rels_path：
    ppt/slides/_rels/slide1.xml.rels -> ppt/slides/slide1.xml
    _rels/.rels -> ""（包级）。
    测试独立实现一份，不导入被测模块的私有函数，避免测试和实现共用
    同一个可能同时出错的地方。
    """
    d = posixpath.dirname(rels_name)
    base = posixpath.basename(rels_name)
    name = base[: -len(".rels")]
    parent_dir = posixpath.dirname(d)
    return posixpath.join(parent_dir, name) if name else parent_dir


def _resolve_target(owner_part: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(posixpath.dirname(owner_part), target))


@pytest.fixture
def deck_with_notes(tmp_path) -> Path:
    """比 deck fixture 多一步：每页加讲稿备注，生成 ppt/notesSlides/*
    part 及其对应的 slide -> notesSlide Relationship（Type 在
    DROP_REL_TYPES 里，不进 keep_parts）。python-pptx 默认模板本身还带
    docProps/thumbnail.jpeg，由包级 _rels/.rels 通过 thumbnail 类型的
    Relationship 引用，同样在 DROP_REL_TYPES 里——两者一起覆盖「part 级
    rels」和「包级 rels」两种悬空 Relationship 场景。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
        s.notes_slide.notes_text_frame.text = f"NOTES-{i + 1}"
    path = tmp_path / "deck_notes.pptx"
    prs.save(path)
    return path


def test_split_has_no_dangling_relationships(deck_with_notes, tmp_path):
    """OPC 规定内部 Relationship 的 Target 必须是包内存在的 part。如果
    只重写了 presentation.xml.rels，而 slide 自己的 .rels（指向已删掉
    的 notesSlide）和包级 _rels/.rels（指向 docProps/thumbnail.jpeg，
    thumbnail 关系在 DROP_REL_TYPES 里、被排除出 keep_parts）照原样搬
    过去，这里就会抓到悬空引用。"""
    out = split_pptx(deck_with_notes, [(1, 3), (4, 8)], tmp_path / "shards")
    for shard in out:
        with zipfile.ZipFile(shard) as z:
            names = set(z.namelist())
            rels_names = [n for n in names if n.endswith(".rels")]
            assert rels_names, f"{shard} 里一个 .rels 都没有，fixture 有问题"
            for rn in rels_names:
                owner = _owner_part(rn)
                root = ET.fromstring(z.read(rn))
                for rel in root:
                    if rel.get("TargetMode") == "External":
                        continue
                    target = _resolve_target(owner, rel.get("Target"))
                    assert target in names, (
                        f"{shard.name} 的 {rn} 里 {rel.get('Id')} "
                        f"指向缺失的 part {target}"
                    )


@pytest.fixture
def deck_with_jump(tmp_path) -> Path:
    """在 slide1 的标题正文里注入一个跳到最后一页的内部超链接
    （<a:hlinkClick r:id="..." action="ppaction://hlinksldjump"/>，
    PowerPoint「链接到幻灯片」/动作按钮生成的真实结构），并在
    slide1.xml.rels 里加一条 Type 与 sldIdLst 里的 slide 关系相同
    （SLIDE_REL_TYPE）的 Relationship 指向最后一页。用来验证内部跳转
    不会把 range 外的整页拖回分片。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
    path = tmp_path / "deck_jump.pptx"
    prs.save(path)

    with zipfile.ZipFile(path) as zin:
        members = {n: zin.read(n) for n in zin.namelist()}

    for prefix, uri in (("a", _A_NS), ("p", _P_NS), ("r", _R_NS)):
        ET.register_namespace(prefix, uri)

    # 从 presentation.xml(.rels) 里找出最后一页的 slide part 文件名，
    # 不硬编码 "slide8.xml"——即便 python-pptx 的创建顺序换了也不受影响。
    pres_root = ET.fromstring(members["ppt/presentation.xml"])
    sld_ids = list(pres_root.find(f"{{{_P_NS}}}sldIdLst"))
    last_rid = sld_ids[-1].get(f"{{{_R_NS}}}id")
    pres_rels_root = ET.fromstring(members["ppt/_rels/presentation.xml.rels"])
    last_slide_target = next(
        rel.get("Target") for rel in pres_rels_root if rel.get("Id") == last_rid
    )
    last_slide_filename = posixpath.basename(last_slide_target)
    assert last_slide_filename != "slide1.xml"

    # 1) slide1.xml.rels 里加一条指向最后一页的 slide 关系
    jump_rid = "rIdJump"
    slide1_rels_name = "ppt/slides/_rels/slide1.xml.rels"
    slide1_rels_root = ET.fromstring(members[slide1_rels_name])
    assert jump_rid not in {rel.get("Id") for rel in slide1_rels_root}
    jump_rel = ET.SubElement(slide1_rels_root, f"{{{_REL_NS}}}Relationship")
    jump_rel.set("Id", jump_rid)
    jump_rel.set("Type", _SLIDE_REL_TYPE)
    jump_rel.set("Target", last_slide_filename)
    members[slide1_rels_name] = ET.tostring(
        slide1_rels_root, xml_declaration=True, encoding="UTF-8"
    )

    # 2) slide1.xml 正文里第一个 run 的 rPr 下挂一个 hlinkClick，指向刚
    #    加的 rId——和真实 PowerPoint「链接到幻灯片」生成的结构一致。
    slide1_root = ET.fromstring(members["ppt/slides/slide1.xml"])
    run = slide1_root.find(f".//{{{_A_NS}}}r")
    assert run is not None
    rpr = run.find(f"{{{_A_NS}}}rPr")
    if rpr is None:
        rpr = ET.Element(f"{{{_A_NS}}}rPr")
        run.insert(0, rpr)
    hlink = ET.SubElement(rpr, f"{{{_A_NS}}}hlinkClick")
    hlink.set(f"{{{_R_NS}}}id", jump_rid)
    hlink.set("action", "ppaction://hlinksldjump")
    members["ppt/slides/slide1.xml"] = ET.tostring(
        slide1_root, xml_declaration=True, encoding="UTF-8"
    )

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in members.items():
            zout.writestr(name, data)

    return path


def test_split_does_not_pull_in_slide_via_internal_hyperlink(deck_with_jump, tmp_path):
    """第 1 页有一个跳到最后一页的内部超链接。切前两页时，这条超链接
    指向的最后一页不应该被拖进分片——分片边界只能由 range 决定，不能
    被 deck 内部的跳转关系改写。"""
    out = split_pptx(deck_with_jump, [(1, 2)], tmp_path / "shards")
    with zipfile.ZipFile(out[0]) as z:
        slide_parts = [n for n in z.namelist() if _SLIDE_PART_RE.match(n)]
    assert len(slide_parts) == 2
    assert probe(out[0]).slide_count == 2


@pytest.fixture
def deck_with_mc_ignorable(deck) -> Path:
    """人工往 presentation.xml 根元素注入真实 PowerPoint 常见的
    xmlns:p14=... mc:Ignorable="p14"。presentation.xml 里没有任何
    p14: 元素，ET 往返只序列化实际用到的命名空间，会把 xmlns:p14
    声明丢掉、只留下指向未声明前缀的 mc:Ignorable="p14"，违反 MCE
    规范——这个 fixture 专门用来暴露这类问题。
    """
    with zipfile.ZipFile(deck) as zin:
        members = {n: zin.read(n) for n in zin.namelist()}

    pres_xml = members["ppt/presentation.xml"].decode("utf-8")
    injected = pres_xml.replace(
        "<p:presentation ",
        '<p:presentation '
        'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
        'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" '
        'mc:Ignorable="p14" ',
        1,
    )
    assert injected != pres_xml, "注入失败，说明 <p:presentation 前缀假设不对"
    members["ppt/presentation.xml"] = injected.encode("utf-8")

    with zipfile.ZipFile(deck, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in members.items():
            zout.writestr(name, data)
    return deck


def test_split_preserves_mc_ignorable_and_its_namespace(deck_with_mc_ignorable, tmp_path):
    out = split_pptx(deck_with_mc_ignorable, [(1, 3)], tmp_path / "shards")
    with zipfile.ZipFile(out[0]) as z:
        pres_xml = z.read("ppt/presentation.xml").decode("utf-8")
    assert (
        'xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main"'
        in pres_xml
    )
    assert 'mc:Ignorable="p14"' in pres_xml
    assert probe(out[0]).slide_count == 3


# ---------------------------------------------------------------------------
# 第二轮代码审查追加：针对 _rewrite_presentation 正则手术本身的三个
# 真实失效模式（均为合法 XML，均端到端跑 split_pptx 复现）：
# A) r:id="rId2" 两侧带空格；B) 关系命名空间绑到非 "r:" 前缀；
# C) <p:sldId> 非自闭合、挂 <p:extLst/> 子元素。
# 详见 task-3-report.md「第二轮代码审查追加修复」一节。
# ---------------------------------------------------------------------------


def _rebuild_zip_with_presentation_xml(src: Path, new_pres_xml: str) -> None:
    """原地重写 src 里的 ppt/presentation.xml，其余 part 原样保留。"""
    with zipfile.ZipFile(src) as zin:
        members = {n: zin.read(n) for n in zin.namelist()}
    members["ppt/presentation.xml"] = new_pres_xml.encode("utf-8")
    with zipfile.ZipFile(src, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in members.items():
            zout.writestr(name, data)


@pytest.fixture
def deck_with_spaced_rid(tmp_path) -> Path:
    """r:id="rIdN" 属性等号两侧带空格——XML 规范里 Eq ::= S? '=' S?，
    这是合法写法。旧的 _RID_ATTR_RE 按 r:id="..." 硬匹配、找不到就走
    「删除」兜底，会把 sldIdLst 整个删空，python-pptx 打开后是 0 页，
    而且没有任何异常，纯粹静默失败。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
    path = tmp_path / "deck_spaced_rid.pptx"
    prs.save(path)

    with zipfile.ZipFile(path) as zin:
        pres_xml = zin.read("ppt/presentation.xml").decode("utf-8")
    injected = pres_xml.replace('r:id="', 'r:id = "')
    assert injected != pres_xml, "注入失败，说明 r:id=\" 的假设不对"
    _rebuild_zip_with_presentation_xml(path, injected)
    return path


def test_split_handles_spaced_rid_attribute(deck_with_spaced_rid, tmp_path):
    out = split_pptx(deck_with_spaced_rid, [(1, 3), (4, 8)], tmp_path / "shards")
    assert probe(out[0]).slide_count == 3
    assert probe(out[1]).slide_count == 5
    assert len(Presentation(str(out[0])).slides) == 3
    assert len(Presentation(str(out[1])).slides) == 5


@pytest.fixture
def deck_with_rel_prefix(tmp_path) -> Path:
    """把 relationships 命名空间的前缀从惯例的 "r" 换成 "rel"——前缀
    本身没有语义，命名空间靠 xmlns:rel="...relationships" 的 URI 绑定，
    这是合法 XML，ElementTree（_slide_order/_read_rels）会正确解析。
    但正则手术是按字面量 "r:id" 匹配的，遇到这种文件应该响亮失败，
    不能把 sldIdLst 删空、静默产出 0 页 deck。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
    path = tmp_path / "deck_rel_prefix.pptx"
    prs.save(path)

    with zipfile.ZipFile(path) as zin:
        pres_xml = zin.read("ppt/presentation.xml").decode("utf-8")
    injected = pres_xml.replace("xmlns:r=", "xmlns:rel=").replace("r:id=", "rel:id=")
    assert injected != pres_xml, "注入失败，说明 xmlns:r=/r:id= 的假设不对"
    assert "rel:id=" in injected
    _rebuild_zip_with_presentation_xml(path, injected)
    return path


def test_split_raises_on_unrecognized_rid_prefix(deck_with_rel_prefix, tmp_path):
    """响亮失败好过静默产出一份 0 页的空 deck。"""
    with pytest.raises(ValueError):
        split_pptx(deck_with_rel_prefix, [(1, 3), (4, 8)], tmp_path / "shards")


@pytest.fixture
def deck_with_non_self_closing_sldid(tmp_path) -> Path:
    """<p:sldId> 挂 <p:extLst/> 子元素时不是自闭合标签（ECMA-376 的
    CT_SlideIdListEntry 允许 extLst），这是合法但少见的形式。只认
    自闭合形式的正则会完全匹配不上，一条都删不掉，留下指向已删关系
    的孤儿 rId——python-pptx 打开分片时会直接抛
    KeyError: "no relationship with key 'rIdN'"。
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
    path = tmp_path / "deck_extlst.pptx"
    prs.save(path)

    with zipfile.ZipFile(path) as zin:
        pres_xml = zin.read("ppt/presentation.xml").decode("utf-8")
    injected = re.sub(
        r"<p:sldId([^>]*?)/>",
        r"<p:sldId\1><p:extLst/></p:sldId>",
        pres_xml,
    )
    assert injected != pres_xml, "注入失败，说明 <p:sldId .../> 的假设不对"
    assert "<p:extLst/></p:sldId>" in injected
    _rebuild_zip_with_presentation_xml(path, injected)
    return path


def test_split_handles_non_self_closing_sldid(deck_with_non_self_closing_sldid, tmp_path):
    out = split_pptx(
        deck_with_non_self_closing_sldid, [(1, 3), (4, 8)], tmp_path / "shards"
    )
    assert probe(out[0]).slide_count == 3
    assert probe(out[1]).slide_count == 5
    # 无孤儿 rId：能被 python-pptx 打开且页数正确。若正则没匹配上、
    # 留下指向已删关系的孤儿 rId，这里会抛
    # KeyError: "no relationship with key 'rIdN'"。
    assert len(Presentation(str(out[0])).slides) == 3
    assert len(Presentation(str(out[1])).slides) == 5
