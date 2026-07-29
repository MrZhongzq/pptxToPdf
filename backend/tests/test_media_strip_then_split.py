"""终审 M-2：媒体剥离 + 分片切片的组合路径。

三期的分片切片（pptx_split.py）和五期的媒体剥离（media_strip.py）都重写
同一份 OPC 包，Task 1 把六个共用工具函数提取到 opc_rewrite.py。生产路径
上（pipeline.run_task）剥离先跑、就地覆盖原件，之后如果还需要切片，
prepare_shards 切的就是"剥离后的文件"——.rels 和 [Content_Types].xml 被
改写了两次。

仓库里这条组合路径原本一条测试都没有：test_media_strip.py 只剥不切；
test_pptx_split.py 的输入是没被剥离过的 deck；test_pipeline_media_strip.py
四条全部 monkeypatch 掉了 strip_media。这里补一条集成测试钉住它，另外
钉住一个本期新引入的不对称：media_strip 的复制循环用负向过滤（只排除
`drop` 里的 part，其余原样搬过去，包括显式目录条目），pptx_split 的复制
循环用正向过滤（只搬 keep_parts 里认得的 part，显式目录条目天然被排除）
——zip 里带显式目录条目时，剥离产出的文件里会多留一条 pptx_split 从不
会写的目录记录，这条记录会作为下一步 split_pptx 的输入，两种过滤写法要
在这个交接点上继续保持互不冲突。
"""

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.media_strip import strip_media
from app.services.opc_rewrite import owner_part, read_rels
from app.services.pptx_split import split_pptx

VIDEO_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/video"
MEDIA_REL = "http://schemas.microsoft.com/office/2007/relationships/media"


def _six_slide_deck(path: Path) -> Path:
    prs = Presentation()
    for i in range(6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        box.text_frame.text = f"PAGE-{i + 1}"
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
    prs.save(str(path))
    return path


def _inject_fake_video(
    src: Path, dest: Path, *, payload_mb: int = 1, with_dir_entry: bool = False
) -> Path:
    """往 slide1 的 rels 塞一段假视频（video + media 两条关系指向同一
    part，与真实 PowerPoint 产出一致），可选再塞一条显式的 zip 目录条目
    ——真实 zip 工具（包括某些版本的 PowerPoint/Office）偶尔会写这种
    记录，media_strip 与 pptx_split 对它的处理方式不同，这里造出来才能
    测到那条不对称。
    """
    payload = b"\x00" * (payload_mb * 1024 * 1024)
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dest, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            raw = zin.read(item.filename)
            if item.filename == "ppt/slides/_rels/slide1.xml.rels":
                injected = (
                    f'<Relationship Id="rIdVid" Type="{VIDEO_REL}" Target="../media/movie1.mp4"/>'
                    f'<Relationship Id="rIdMed" Type="{MEDIA_REL}" Target="../media/movie1.mp4"/>'
                ).encode()
                raw = raw.replace(b"</Relationships>", injected + b"</Relationships>")
            elif item.filename == "[Content_Types].xml":
                raw = raw.replace(
                    b"</Types>",
                    b'<Default Extension="mp4" ContentType="video/mp4"/></Types>',
                )
            zout.writestr(item, raw)
        zout.writestr("ppt/media/movie1.mp4", payload)
        if with_dir_entry:
            # 显式目录条目：名字以 "/" 结尾、无内容——media_strip 的复制
            # 循环只排除 `drop` 里的 part，这条记录不在 drop 里，会被原样
            # 搬进剥离后的文件；pptx_split 的复制循环只搬 keep_parts 认得
            # 的 part，这条记录天然不在 keep_parts 里，会被丢弃。
            zout.writestr(zipfile.ZipInfo("ppt/media/"), b"")
    return dest


def _assert_no_dangling_rels(path: Path) -> None:
    with zipfile.ZipFile(path) as zf:
        names = set(zf.namelist())
        for name in names:
            if not name.endswith(".rels"):
                continue
            owner = owner_part(name)
            for _rid, _type, target in read_rels(zf, owner):
                assert target in names, f"{path.name}:{name} 指向不存在的 {target}"


def _strip_then_split(tmp_path: Path, *, with_dir_entry: bool) -> list[Path]:
    base = _six_slide_deck(tmp_path / "base.pptx")
    deck = tmp_path / "deck.pptx"
    _inject_fake_video(base, deck, with_dir_entry=with_dir_entry)

    strip_result = strip_media(deck)
    assert strip_result.stripped is True
    assert "ppt/media/movie1.mp4" not in {
        n for n in zipfile.ZipFile(deck).namelist()
    }

    out_dir = tmp_path / "shards"
    return split_pptx(deck, [(1, 2), (3, 4), (5, 6)], out_dir)


def test_strip_then_split_produces_three_valid_shards(tmp_path):
    """真实组合：6 页含 1MB 假视频的 deck -> 真 strip_media -> 真
    split_pptx([(1,2),(3,4),(5,6)])。三个分片都必须 .rels 无悬空、
    python-pptx 能打开、页数正确、且不含已被剥离的视频 part。
    """
    shards = _strip_then_split(tmp_path, with_dir_entry=False)

    assert len(shards) == 3
    for shard in shards:
        _assert_no_dangling_rels(shard)
        prs = Presentation(str(shard))
        assert len(prs.slides) == 2
        with zipfile.ZipFile(shard) as zf:
            assert "ppt/media/movie1.mp4" not in zf.namelist()


def test_strip_then_split_with_explicit_zip_directory_entry(tmp_path):
    """变体：zip 里带显式目录条目。media_strip 的复制循环是负向过滤
    （排除 drop，其余原样搬），pptx_split 的复制循环是正向过滤（只搬
    keep_parts 认得的 part）——这条不对称是本期新引入的，剥离产出的文件
    里会多留一条 pptx_split 从不会写的目录记录，下一步切片必须仍然正常
    工作，不能因为多出这条记录而炸掉或产出损坏的分片。
    """
    shards = _strip_then_split(tmp_path, with_dir_entry=True)

    assert len(shards) == 3
    for shard in shards:
        _assert_no_dangling_rels(shard)
        prs = Presentation(str(shard))
        assert len(prs.slides) == 2
