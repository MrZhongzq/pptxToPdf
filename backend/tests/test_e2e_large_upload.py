import pytest
from httpx import ASGITransport, AsyncClient
from pptx import Presentation
from pptx.util import Emu
from pypdf import PdfReader

from app.config import settings
from app.db import Base, engine

SLIDES = 120


@pytest.fixture
def client(tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "storage_root", tmp_path / "storage")
    # 简报原值 256 * 1024：120 张空白版式幻灯片实测仅约 110 KiB，
    # 256 KiB 块下只会切出 1 块，"assert total > 1" 必然失败，
    # 中断续传场景根本没被覆盖到。调小到 16 KiB 以确保样本稳定切出
    # 多块（本机实测 7 块），真正覆盖"传一半、查状态、续传剩余"链路。
    monkeypatch.setattr(settings, "chunk_size", 16 * 1024)
    settings.ensure_dirs()
    Base.metadata.drop_all(engine)
    Base.metadata.create_all(engine)

    from app.main import app

    return AsyncClient(transport=ASGITransport(app=app), base_url="http://test")


@pytest.fixture
def big_pptx(tmp_path):
    """120 页 deck —— 正是 Graph 转不了、必须走 LibreOffice 的那一类。"""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    for _ in range(SLIDES):
        prs.slides.add_slide(prs.slide_layouts[6])
    path = tmp_path / "big.pptx"
    prs.save(path)
    return path.read_bytes()


async def test_interrupted_upload_resumes_and_completes(client, big_pptx, tmp_path):
    created = (
        await client.post(
            "/api/uploads", json={"filename": "big.pptx", "size": len(big_pptx)}
        )
    ).json()
    uid, size, total = (
        created["upload_id"],
        created["chunk_size"],
        created["total_chunks"],
    )
    assert total > 1, "样本需要多于一块才能验证续传"

    # 只传前半段，模拟中途断线
    half = total // 2
    for idx in range(half):
        await client.put(
            f"/api/uploads/{uid}/chunks/{idx}",
            content=big_pptx[idx * size : (idx + 1) * size],
        )

    status = (await client.get(f"/api/uploads/{uid}")).json()
    assert status["received_indices"] == list(range(half))

    # 续传剩余块
    for idx in range(half, total):
        await client.put(
            f"/api/uploads/{uid}/chunks/{idx}",
            content=big_pptx[idx * size : (idx + 1) * size],
        )

    task_id = (await client.post(f"/api/uploads/{uid}/complete")).json()["task_id"]

    task = (await client.get(f"/api/tasks/{task_id}")).json()
    assert task["status"] == "done"
    assert task["slide_count"] == SLIDES

    pdf = tmp_path / "out.pdf"
    pdf.write_bytes((await client.get(f"/api/tasks/{task_id}/download")).content)
    assert len(PdfReader(pdf).pages) == SLIDES
