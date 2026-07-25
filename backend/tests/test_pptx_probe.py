import zipfile

import pytest
from pptx import Presentation
from pptx.util import Emu

from app.errors import PptxEncrypted, PptxInvalidZip, PptxNotPresentation
from app.services.pptx_probe import probe

CFB_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


@pytest.fixture
def sample_pptx(tmp_path):
    """用 python-pptx 造样本——仅测试用，生产代码不得依赖它。"""
    prs = Presentation()
    prs.slide_width = Emu(12192000)   # 16:9
    prs.slide_height = Emu(6858000)
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"第 {i + 1} 页"
    path = tmp_path / "sample.pptx"
    prs.save(path)
    return path


def test_probe_slide_count_and_size(sample_pptx):
    meta = probe(sample_pptx)

    assert meta.slide_count == 3
    assert meta.slide_width_emu == 12192000
    assert meta.slide_height_emu == 6858000


def test_probe_collects_fonts(sample_pptx):
    meta = probe(sample_pptx)

    assert isinstance(meta.fonts, tuple)
    assert all(isinstance(f, str) for f in meta.fonts)
    assert meta.fonts == tuple(sorted(meta.fonts))


def test_probe_rejects_non_zip(tmp_path):
    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"this is not a zip file at all")

    with pytest.raises(PptxInvalidZip):
        probe(bad)


def test_probe_rejects_encrypted(tmp_path):
    enc = tmp_path / "enc.pptx"
    enc.write_bytes(CFB_MAGIC + b"\x00" * 128)

    with pytest.raises(PptxEncrypted):
        probe(enc)


def test_probe_rejects_zip_without_presentation(tmp_path):
    zpath = tmp_path / "notppt.pptx"
    with zipfile.ZipFile(zpath, "w") as zf:
        zf.writestr("hello.txt", "world")

    with pytest.raises(PptxNotPresentation):
        probe(zpath)
