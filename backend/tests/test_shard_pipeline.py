"""分片流水线的纯逻辑测试。

不需要 Azure 凭证：引擎一律用假引擎注入，覆盖的是编排本身——页范围规划、
切片后的体积实测与重切、分片数/合并体积上限、状态机转换、合并顺序、
失败传播、中间产物清理。

测试分三层，缺一层就会留下"改一行悄悄退回"的缺口：
1. 子函数层（prepare_shards / convert_shard / merge_shards 各自的行为）；
2. 接线层（run_task 到底有没有走分片分支、有没有把 shard_ids 交给队列）；
3. 队列层（enqueue_shards 到底有没有用 Dependency(allow_failure=True)）。
"""
import io
import random
import shutil
from datetime import datetime, timedelta, timezone
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches
from pypdf import PdfReader
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

import app.services.pipeline as pipeline_module
import app.services.retention as retention_module
import app.services.shard_pipeline as shard_module
from app.config import settings
from app.errors import ConversionTimeout, GraphNotConfigured, ShardBudgetExceeded, ShardTooLarge
from app.models import Task, TaskShard
from app.services.engine_router import select_engine as real_select_engine
from app.services.graph_credentials import save_credentials
from app.services.pipeline import compute_timeout_s
from app.services.retention import reap_stale_tasks
from app.services.shard_pipeline import (
    convert_shard,
    merge_shards,
    prepare_shards,
    shard_dir,
)

MIB = 1024 * 1024

DECK_SLIDES = 8

# 审查轮之后（I2）：run_task 进分片分支前会真的查一次 Graph 是否配置
# （is_graph_configured(session)），这行检查独立于 get_engine 有没有被
# fake 掉——所以 wired_pipeline 系的测试即便全程用假引擎，也得先在凭证表
# 里放一行能通过 load_credentials 的记录，否则会在 prepare_shards 之前
# 就被新加的早退检查拦下。这里用固定的假 Fernet key，凭证内容本身没人
# 校验，只要求解密成功。
GRAPH_TEST_SECRET_KEY = "8I3F3CqPwlEsmMDLbEIVSXd8oXlmqkOMWFnDPbNXKvA="


# ---------------------------------------------------------------- 素材构造


def _pdf(path: Path, labels: list[str]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    c = canvas.Canvas(str(path), pagesize=letter)
    for label in labels:
        c.setFont("Helvetica", 40)
        c.drawCentredString(300, 400, label)
        c.showPage()
    c.save()
    return path


HEAVY_SLIDES = 2
"""前两页各嵌一张约 1MB 的噪声图，其余页只有标题。

刻意做成极端失真的体积分布——这正是 plan_ranges 把终局判定权交给调用方的
理由。若每页体积均匀，「均摊估算」和「实测」就几乎重合，重切逻辑永远触发
不到，这个文件里关于实测复核的测试也就测不到东西。
"""


def _noisy_png(seed: int, w: int = 700, h: int = 450) -> io.BytesIO:
    """随机噪声图几乎不可压缩，单张约 1MB。

    seed 必须逐页不同：python-pptx 按图片字节做去重，两页塞同一张图会
    合并成同一个 media part，体积分布就又变均匀了。
    """
    rng = random.Random(seed)
    img = Image.new("RGB", (w, h))
    img.putdata(
        [
            (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
            for _ in range(w * h)
        ]
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


@pytest.fixture(scope="session")
def _deck_master(tmp_path_factory) -> Path:
    """构造一次、全模块复用——噪声图生成是这个文件里最慢的一步。"""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(12192000), Emu(6858000)
    for i in range(DECK_SLIDES):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"PAGE-{i + 1}"
        if i < HEAVY_SLIDES:
            s.shapes.add_picture(_noisy_png(seed=i), Inches(1), Inches(2), width=Inches(4))
    path = tmp_path_factory.mktemp("deckmaster") / "deck.pptx"
    prs.save(path)
    return path


@pytest.fixture
def deck(tmp_path, _deck_master) -> Path:
    dst = tmp_path / "deck.pptx"
    shutil.copyfile(_deck_master, dst)
    return dst


@pytest.fixture
def tiny_deck_factory(tmp_path):
    """无图片的小 deck，用于只关心页数的分片转换测试。"""

    def _make(path: Path, pages: int) -> Path:
        prs = Presentation()
        for i in range(pages):
            s = prs.slides.add_slide(prs.slide_layouts[5])
            s.shapes.title.text = f"T{i + 1}"
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        return path

    return _make


@pytest.fixture
def storage(tmp_path, monkeypatch):
    monkeypatch.setattr(settings, "storage_root", tmp_path / "storage")
    settings.ensure_dirs()
    return settings.storage_root


class _FakeEngine:
    """假转换引擎：记录每次调用的参数，按需产出 PDF 或抛错。"""

    def __init__(self, pages: int | None = None, exc: Exception | None = None) -> None:
        self.calls: list[dict] = []
        self.pages = pages
        self.exc = exc

    def convert(self, src: Path, meta, dest: Path, *, timeout_s: float) -> None:
        self.calls.append(
            {
                "src": Path(src),
                "dest": Path(dest),
                "timeout_s": timeout_s,
                "slide_count": meta.slide_count,
            }
        )
        if self.exc is not None:
            raise self.exc
        n = self.pages if self.pages is not None else meta.slide_count
        _pdf(Path(dest), [f"X{i}" for i in range(n)])


def _install_engine(monkeypatch, engine: _FakeEngine, target=shard_module) -> list[str]:
    """接管 get_engine，并把请求过的引擎名记下来——"有没有偷偷换引擎"
    是本期的红线，必须可断言。

    `**_kwargs` 吞掉 `session=`：真实 get_engine 现在需要 session 才能给
    graph 引擎注入凭证（Task 8），但这里假引擎完全绕过了凭证加载这条路径，
    调用方仍然会传 session= 过来，必须接住而不是 TypeError。
    """
    asked: list[str] = []

    def _fake_get_engine(name: str, **_kwargs):
        asked.append(name)
        return engine

    monkeypatch.setattr(target, "get_engine", _fake_get_engine)
    return asked


@pytest.fixture
def sharded_task(session, storage):
    task = Task(
        task_id="T1",
        upload_id="U1",
        original_filename="deck.pptx",
        size_bytes=100,
        slide_count=3,
        engine="graph",
        status="converting",
        shard_total=2,
    )
    session.add(task)
    # 故意倒序插入：SQLite 的自然 rowid 顺序会与 index 顺序相反，于是
    # merge_shards 里那条 `.order_by(TaskShard.index)` 一旦被删掉，页序
    # 立刻错乱。顺序插入时 rowid 顺序恰好正确，这条防线就没人守。
    for i, (ps, pe) in reversed(list(enumerate([(1, 2), (3, 3)]))):
        session.add(
            TaskShard(
                shard_id=f"S{i}",
                task_id="T1",
                index=i,
                page_start=ps,
                page_end=pe,
                status="pending",
            )
        )
    session.commit()
    return task


@pytest.fixture
def use_test_session(session, monkeypatch):
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    return session


# ---------------------------------------------------------------- merge_shards


def test_merge_marks_done_when_all_shards_succeed(
    session, sharded_task, use_test_session
):
    d = shard_dir("T1")
    for i, labels in enumerate([["P1", "P2"], ["P3"]]):
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status = "done"
        shard.output_path = str(p)
    session.commit()

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "done"
    assert task.output_path is not None
    assert len(PdfReader(task.output_path).pages) == 3


def test_merge_fails_when_any_shard_failed(session, sharded_task, use_test_session):
    """9 成 1 败也必须整体失败——合并出一份缺了中间几页的 PDF，
    比明确报错糟糕得多。"""
    d = shard_dir("T1")
    p = _pdf(d / "000.pdf", ["P1", "P2"])
    s0 = session.get(TaskShard, "S0")
    s0.status, s0.output_path = "done", str(p)
    s1 = session.get(TaskShard, "S1")
    s1.status, s1.error_code, s1.error_message = "failed", "CONVERSION_TIMEOUT", "转换超时"
    session.commit()

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "failed"
    assert task.error_code == "CONVERSION_TIMEOUT"
    assert task.output_path is None


def test_merge_fails_when_shards_unfinished(session, sharded_task, use_test_session):
    """没有失败但也没跑完——汇总 job 被提前触发时不能当成功处理。"""
    d = shard_dir("T1")
    p = _pdf(d / "000.pdf", ["P1", "P2"])
    s0 = session.get(TaskShard, "S0")
    s0.status, s0.output_path = "done", str(p)
    session.commit()  # S1 仍是 pending

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "failed"
    assert task.output_path is None


def test_merge_fails_when_page_count_mismatches(session, sharded_task, use_test_session):
    """合并结果的总页数必须等于 slide_count。"""
    d = shard_dir("T1")
    for i, labels in enumerate([["P1"], ["P3"]]):  # 少了一页
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status, shard.output_path = "done", str(p)
    session.commit()

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "failed"
    assert task.error_code == "CONVERSION_PAGE_MISMATCH"
    assert task.output_path is None


def test_merge_uses_index_order_not_filename_order(
    session, sharded_task, use_test_session
):
    """合并顺序只由 TaskShard.index 决定。文件名故意反着排：任何形式的
    按路径排序都会让页序变成 P3,P1,P2。"""
    d = shard_dir("T1")
    p0 = _pdf(d / "zzz.pdf", ["P1", "P2"])
    p1 = _pdf(d / "aaa.pdf", ["P3"])
    for sid, p in (("S0", p0), ("S1", p1)):
        shard = session.get(TaskShard, sid)
        shard.status, shard.output_path = "done", str(p)
    session.commit()

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "done"
    texts = [page.extract_text().strip() for page in PdfReader(task.output_path).pages]
    assert texts == ["P1", "P2", "P3"]


def test_merge_rejects_when_total_bytes_exceed_budget(
    session, sharded_task, use_test_session, monkeypatch
):
    """merge_pdfs 会把所有分片一次性载入同一个 PdfWriter（pypdf 没有流式
    合并），峰值约 2.2 倍输入体积。超预算必须在载入之前明确报错，
    而不是让 worker 被 OOM killer 静默干掉、任务永远卡在 running。"""
    d = shard_dir("T1")
    for i, labels in enumerate([["P1", "P2"], ["P3"]]):
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status, shard.output_path = "done", str(p)
    session.commit()
    monkeypatch.setattr(settings, "graph_max_merge_bytes", 10)

    called: list = []
    monkeypatch.setattr(
        shard_module, "merge_pdfs", lambda *a, **k: called.append(a) or 0
    )

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "failed"
    assert task.error_code == "SHARD_BUDGET_EXCEEDED"
    assert called == []  # 关键：根本没进 merge_pdfs


def test_merge_cleans_shard_dir(session, sharded_task, use_test_session):
    d = shard_dir("T1")
    for i, labels in enumerate([["P1", "P2"], ["P3"]]):
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status, shard.output_path = "done", str(p)
    session.commit()

    merge_shards("T1")
    # 中间产物体积是原文件两倍（分片 pptx + 分片 PDF），必须清掉
    assert not d.exists()


def test_merge_cleans_shard_dir_on_failure(session, sharded_task, use_test_session):
    d = shard_dir("T1")
    d.mkdir(parents=True, exist_ok=True)
    (d / "000.pptx").write_bytes(b"leftover")
    s1 = session.get(TaskShard, "S1")
    s1.status, s1.error_code = "failed", "CONVERSION_FAILED"
    session.commit()

    merge_shards("T1")
    assert not d.exists()


def test_merge_ignores_missing_task(use_test_session):
    merge_shards("does-not-exist")  # 不许抛


def test_merge_does_not_revive_a_task_already_in_terminal_state(
    session, sharded_task, use_test_session
):
    """终态翻转比漏合并更糟。

    孤儿回收器把任务标 failed（前端一见 failed 就停止轮询、原文件已删）之后，
    迟到的分片仍可能陆续跑完并触发汇总——若这里不设闸门，任务会从 failed
    一路被改回 done，而没有任何人还在看它。
    """
    d = shard_dir("T1")
    for i, labels in enumerate([["P1", "P2"], ["P3"]]):
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status, shard.output_path = "done", str(p)
    task = session.get(Task, "T1")
    task.status, task.error_code = "failed", "TASK_ABANDONED"
    session.commit()

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "failed"
    assert task.error_code == "TASK_ABANDONED"
    assert task.output_path is None


def test_merge_is_idempotent(session, sharded_task, use_test_session):
    """重跑不许把已经 done 的任务降级成 failed。

    第一次的 finally 已经 rmtree 掉分片目录，第二次 _collect_parts 必然在
    「输出文件已丢失」上炸——而那份 PDF 其实完好无损。RQ 的 job 重试、
    人工重放都会踩到这条。
    """
    d = shard_dir("T1")
    for i, labels in enumerate([["P1", "P2"], ["P3"]]):
        p = _pdf(d / f"{i:03d}.pdf", labels)
        shard = session.get(TaskShard, f"S{i}")
        shard.status, shard.output_path = "done", str(p)
    session.commit()

    merge_shards("T1")
    first = session.get(Task, "T1").output_path
    assert first is not None

    merge_shards("T1")

    task = session.get(Task, "T1")
    assert task.status == "done"
    assert task.output_path == first


# ---------------------------------------------------- 孤儿回收与分片的活性信号


def _stale(minutes: int) -> datetime:
    """naive UTC——SQLite dialect 落库时会丢时区，retention 也用 naive 比较。"""
    return datetime.now(timezone.utc).replace(tzinfo=None) - timedelta(minutes=minutes)


@pytest.fixture
def use_retention_session(session, monkeypatch):
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    return session


def _sharded_task_with_shards(
    session, *, task_updated_min: int, shard_updates: list[int], status: str = "converting"
) -> Task:
    task = Task(
        task_id="TR",
        upload_id="UR",
        original_filename="deck.pptx",
        size_bytes=100,
        slide_count=8,
        engine="graph",
        status=status,
        shard_total=len(shard_updates),
        updated_at=_stale(task_updated_min),
    )
    session.add(task)
    for i, minutes in enumerate(shard_updates):
        session.add(
            TaskShard(
                shard_id=f"R{i}",
                task_id="TR",
                index=i,
                page_start=i + 1,
                page_end=i + 1,
                status="done",
                updated_at=_stale(minutes),
            )
        )
    session.commit()
    return task


def test_reap_spares_sharded_task_whose_shards_are_still_progressing(
    session, use_retention_session
):
    """分片路径打破了二期的不变量：单次转换的墙钟被 convert_timeout_max_s
    (1800s) 封顶 < 45 分钟，而分片总墙钟是 N 片之和，12 × 1800s = 6 小时。

    Task 行从 prepare_shards 之后就不再被写（子 job 只写自己那行 TaskShard，
    这是对的、不能改），所以 Task.updated_at 不再是活性信号——拿它判 stale
    会把正在正常转换的 500MB 课件误杀。
    """
    _sharded_task_with_shards(session, task_updated_min=46, shard_updates=[40, 30, 0])

    assert reap_stale_tasks() == 0

    task = session.get(Task, "TR")
    assert task.status == "converting"
    assert task.error_code is None


def test_reap_still_kills_sharded_task_when_shards_also_went_quiet(
    session, use_retention_session
):
    """闸门不能开成永不回收：分片自己也超过 45 分钟没动静，说明 worker
    是真的死了，仍要回收。"""
    _sharded_task_with_shards(session, task_updated_min=60, shard_updates=[50, 47, 46])

    assert reap_stale_tasks() == 1

    task = session.get(Task, "TR")
    assert task.status == "failed"
    assert task.error_code == "TASK_ABANDONED"


def test_reap_covers_merging_state(session, use_retention_session):
    """merge job 被部署/OOM/kill 打断后任务停在 merging。它不在
    NON_TERMINAL 里的话回收器捞不到，任务永远卡住、前端只能轮询到超时——
    正是「分片上限」那条硬要求的立论前提要消灭的失败模式。"""
    from app.services.retention import NON_TERMINAL

    assert "merging" in NON_TERMINAL

    _sharded_task_with_shards(
        session, task_updated_min=60, shard_updates=[50, 50], status="merging"
    )

    assert reap_stale_tasks() == 1
    assert session.get(Task, "TR").status == "failed"


def test_reap_leaves_unsharded_tasks_on_the_original_rule(
    session, use_retention_session
):
    """二期原路径（shard_total 为 None）的判据一个字不改。"""
    session.add(
        Task(
            task_id="TS",
            upload_id="US",
            original_filename="deck.pptx",
            size_bytes=100,
            engine="libreoffice",
            status="converting",
            updated_at=_stale(46),
        )
    )
    session.commit()

    assert reap_stale_tasks() == 1
    assert session.get(Task, "TS").status == "failed"


# ---------------------------------------------------------------- 预算不变量


def test_merge_budget_leaves_headroom_on_a_2gb_worker():
    """上限不是随手定的数，而是由实测倍率反推的：审查实测 4 片 54.1MB 的
    图片密集型 PDF 峰值 162.9MB = 3.01×（tracemalloc，还不含解释器基线与
    分配器碎片）。上限 × 倍率必须明显低于 2GB worker 的可用内存。"""
    peak_bytes = settings.graph_max_merge_bytes * 3.01
    assert peak_bytes < 1024 * MIB, (
        f"合并峰值外推 {peak_bytes / MIB:.0f}MB，2GB worker 上余量不足"
    )


def test_graph_path_capacity_is_not_larger_than_the_advertised_upload_limit():
    """Graph 路径的实际容量 = graph_max_shards × graph_max_shard_bytes，
    天然低于 max_file_size（Graph 的固有限制，不是配置错误）。这里只钉住
    「不许靠调大片数去对齐 600MB」——那会直接 OOM。"""
    capacity = settings.graph_max_shards * settings.graph_max_shard_bytes
    assert capacity <= settings.max_file_size


# ---------------------------------------------------------------- convert_shard


@pytest.fixture
def one_shard(session, storage, tiny_deck_factory):
    task = Task(
        task_id="T2",
        upload_id="U2",
        original_filename="deck.pptx",
        size_bytes=1000,
        slide_count=4,
        engine="graph",
        status="converting",
        shard_total=1,
    )
    session.add(task)
    session.add(
        TaskShard(
            shard_id="SA",
            task_id="T2",
            index=0,
            page_start=1,
            page_end=2,
            status="pending",
        )
    )
    session.commit()
    tiny_deck_factory(shard_dir("T2") / "000.pptx", 2)
    return task


def test_convert_shard_records_output_and_uses_task_engine(
    session, one_shard, use_test_session, monkeypatch
):
    engine = _FakeEngine()
    asked = _install_engine(monkeypatch, engine)

    convert_shard("SA")

    shard = session.get(TaskShard, "SA")
    assert shard.status == "done"
    assert shard.output_path == str((shard_dir("T2") / "000.pdf").resolve())
    # 绝不静默换引擎：用的必须是 Task.engine
    assert asked == ["graph"]
    call = engine.calls[0]
    assert call["src"] == shard_dir("T2") / "000.pptx"
    assert call["dest"] == shard_dir("T2") / "000.pdf"
    assert call["slide_count"] == 2
    # 预算必须按这一片自己的页数与体积算。brief 原本写的是
    # graph_request_timeout_s(50s)——那是单个 HTTP 请求的超时，拿它当整片的
    # 总墙钟预算会让 40MB 分片上传未完即判超时，而且 GraphEngine 里
    # 「wait >= remaining 就放弃」的逻辑会让 Retry-After: 60~300 的退避重试
    # 一次都发不出去。断言必须钉住这个偏离，否则后来者照 brief 改回去没人拦。
    src_size = (shard_dir("T2") / "000.pptx").stat().st_size
    assert call["timeout_s"] == compute_timeout_s(2, src_size)
    assert call["timeout_s"] > settings.graph_request_timeout_s
    # 分片 job 不许碰主任务行——并发写同一行在 SQLite 上会丢更新
    assert session.get(Task, "T2").status == "converting"


def test_convert_shard_records_app_error(session, one_shard, use_test_session, monkeypatch):
    _install_engine(monkeypatch, _FakeEngine(exc=ConversionTimeout("超时了")))

    convert_shard("SA")

    shard = session.get(TaskShard, "SA")
    assert shard.status == "failed"
    assert shard.error_code == "CONVERSION_TIMEOUT"
    assert shard.error_message == "超时了"
    assert session.get(Task, "T2").status == "converting"


def test_convert_shard_records_internal_error_on_crash(
    session, one_shard, use_test_session, monkeypatch
):
    """裸异常不能让 RQ 把 job 记成崩溃而分片状态停在 converting——
    那样汇总 job 会看到一个永远不落地的分片。"""
    _install_engine(monkeypatch, _FakeEngine(exc=RuntimeError("boom")))

    convert_shard("SA")

    shard = session.get(TaskShard, "SA")
    assert shard.status == "failed"
    assert shard.error_code == "INTERNAL_ERROR"


def test_convert_shard_fails_when_source_missing(
    session, one_shard, use_test_session, monkeypatch
):
    (shard_dir("T2") / "000.pptx").unlink()
    _install_engine(monkeypatch, _FakeEngine())

    convert_shard("SA")

    shard = session.get(TaskShard, "SA")
    assert shard.status == "failed"
    assert shard.error_code is not None


def test_convert_shard_ignores_missing_shard(use_test_session):
    convert_shard("nope")  # 不许抛


def test_convert_shard_fails_loudly_when_graph_not_configured(
    session, one_shard, use_test_session, monkeypatch
):
    """红线覆盖分片路径：每个分片是独立 RQ job，各自调用 get_engine("graph",
    session=...)，凭证检查必须在每个分片里都生效，不能只在非分片路径测过。

    刻意不用 _install_engine——要走真实的 get_engine，验证 convert_shard
    里 `session=session` 这个接线真的传下去了。额外 patch httpx.Client 到
    会炸的桩，证明没有任何 HTTP 请求被发出。"""
    import app.services.engines.graph as graph_module

    monkeypatch.setattr(settings, "secret_key", None)  # 未配置

    def _boom(*a, **kw):
        raise AssertionError("不该发出任何 HTTP 请求——凭证检查必须先失败")

    monkeypatch.setattr(graph_module.httpx, "Client", _boom)

    convert_shard("SA")

    shard = session.get(TaskShard, "SA")
    assert shard.status == "failed"
    assert shard.error_code == "GRAPH_NOT_CONFIGURED"
    # 分片失败不许连带把主任务标终态——那是汇总 job 的职责
    assert session.get(Task, "T2").status == "converting"


# ---------------------------------------------------------------- prepare_shards


def _make_task(session, deck: Path, task_id: str = "T3") -> Task:
    src = settings.originals_dir / f"{task_id}.pptx"
    shutil.copyfile(deck, src)
    task = Task(
        task_id=task_id,
        upload_id="U3",
        original_filename="deck.pptx",
        size_bytes=src.stat().st_size,
        slide_count=DECK_SLIDES,
        engine="graph",
        requested_engine="graph",
        status="converting",
    )
    session.add(task)
    session.commit()
    return task


def _make_fresh_task(
    session, deck: Path, task_id: str, requested_engine: str | None
) -> Task:
    """跟 _make_task 不同：不预置 engine/slide_count，让 run_task 自己走
    probe + select_engine 全流程决定——C1 的红线测试要验证的正是这条
    "requested_engine 字段真的传到了 select_engine 的 requested= 形参"
    的接线，预置好 engine 字段会绕过这条接线，测不出漏传的问题。"""
    src = settings.originals_dir / f"{task_id}.pptx"
    shutil.copyfile(deck, src)
    task = Task(
        task_id=task_id,
        upload_id="U-fresh",
        original_filename="deck.pptx",
        size_bytes=src.stat().st_size,
        requested_engine=requested_engine,
        status="pending",
    )
    session.add(task)
    session.commit()
    return task


def test_prepare_shards_creates_rows_and_files(session, storage, deck, monkeypatch):
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 2)
    task = _make_task(session, deck)
    src = settings.originals_dir / "T3.pptx"

    shard_ids = prepare_shards(session, task, src, src.stat().st_size)

    assert len(shard_ids) == 4
    assert task.shard_total == 4
    shards = (
        session.query(TaskShard)
        .filter(TaskShard.task_id == "T3")
        .order_by(TaskShard.index)
        .all()
    )
    assert [s.index for s in shards] == [0, 1, 2, 3]
    assert [(s.page_start, s.page_end) for s in shards] == [(1, 2), (3, 4), (5, 6), (7, 8)]
    for s in shards:
        assert (shard_dir("T3") / f"{s.index:03d}.pptx").is_file()
    # 中间的 scratch 目录必须清干净，只留最终分片
    assert sorted(p.name for p in shard_dir("T3").iterdir()) == [
        "000.pptx",
        "001.pptx",
        "002.pptx",
        "003.pptx",
    ]


def test_prepare_shards_resplits_when_measured_size_exceeds_budget(
    session, storage, deck, monkeypatch
):
    """plan_ranges 只有「总体积 ÷ 页数」的均值可用，共享的 theme/master 在
    每片里各留一份，所以实测总是比均摊估算大。调用方必须实测复核并再切一轮，
    否则规划器放行的分片会直接怼给 Graph，换回一个难懂的 Graph 错误。"""
    from app.services.shard_planner import plan_ranges

    task = _make_task(session, deck)
    src = settings.originals_dir / "T3.pptx"
    size = src.stat().st_size
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", DECK_SLIDES)
    # 体积几乎全集中在前两页。预算取全包的 3/4：
    # - 均摊估算看到的每页只有 size/8，于是认为 6 页一片都塞得下 → 规划 2 片；
    # - 但真正装下前两页的那一片实测接近整包体积，必然超预算，只有把两张
    #   重图拆到不同分片才收敛。
    monkeypatch.setattr(settings, "graph_max_shard_bytes", int(size * 0.75))

    planned = plan_ranges(
        DECK_SLIDES, size, settings.graph_max_pages_per_shard, settings.graph_max_shard_bytes
    )
    assert len(planned) == 2  # 估算说 2 片就够

    shard_ids = prepare_shards(session, task, src, size)

    assert len(shard_ids) > len(planned)  # 实测复核之后必须切得更细
    shards = (
        session.query(TaskShard)
        .filter(TaskShard.task_id == "T3")
        .order_by(TaskShard.index)
        .all()
    )
    # 覆盖完整、无缝无重叠
    assert shards[0].page_start == 1
    assert shards[-1].page_end == DECK_SLIDES
    for prev, nxt in zip(shards, shards[1:]):
        assert nxt.page_start == prev.page_end + 1
    # 每一片的实测体积都在预算内——这才是本条契约的终局判定
    for s in shards:
        path = shard_dir("T3") / f"{s.index:03d}.pptx"
        assert path.stat().st_size <= settings.graph_max_shard_bytes


def test_prepare_shards_raises_when_single_page_too_large(
    session, storage, deck, monkeypatch
):
    """单片已经只剩 1 页却仍然超限，没有再切的余地——必须响亮失败，
    不能无限重切，也不能静默放行。"""
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", DECK_SLIDES)
    monkeypatch.setattr(settings, "graph_max_shard_bytes", 1)
    monkeypatch.setattr(settings, "graph_max_shards", 1000)
    task = _make_task(session, deck)
    src = settings.originals_dir / "T3.pptx"

    with pytest.raises(ShardTooLarge):
        prepare_shards(session, task, src, src.stat().st_size)

    assert task.shard_total is None
    assert not shard_dir("T3").exists()  # 失败路径也要清理中间产物


def test_prepare_shards_rejects_too_many_shards_at_planning(
    session, storage, deck, monkeypatch
):
    """分片总数必须有显式上限：merge_pdfs 的峰值内存正比于分片总量。
    规划阶段就已经超限时快速失败，不必先切一整轮 IO 再拒绝。"""
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 1)
    monkeypatch.setattr(settings, "graph_max_shards", 3)
    task = _make_task(session, deck)
    src = settings.originals_dir / "T3.pptx"

    with pytest.raises(ShardBudgetExceeded):
        prepare_shards(session, task, src, src.stat().st_size)

    assert session.query(TaskShard).filter(TaskShard.task_id == "T3").count() == 0
    assert not shard_dir("T3").exists()


def test_prepare_shards_rejects_when_resplit_exceeds_shard_cap(
    session, storage, deck, monkeypatch
):
    """上限必须在实测复核之后再判一次。重切会让分片数比规划值多，只卡规划
    输出等于没卡——这里规划只要 2 片（远在上限内），重切之后才超。"""
    from app.services.shard_planner import plan_ranges

    task = _make_task(session, deck)
    src = settings.originals_dir / "T3.pptx"
    size = src.stat().st_size
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", DECK_SLIDES)
    monkeypatch.setattr(settings, "graph_max_shard_bytes", int(size * 0.75))
    monkeypatch.setattr(settings, "graph_max_shards", 3)

    assert (
        len(
            plan_ranges(
                DECK_SLIDES,
                size,
                settings.graph_max_pages_per_shard,
                settings.graph_max_shard_bytes,
            )
        )
        <= settings.graph_max_shards
    )  # 规划阶段的快速失败这次不会触发

    with pytest.raises(ShardBudgetExceeded):
        prepare_shards(session, task, src, size)

    assert session.query(TaskShard).filter(TaskShard.task_id == "T3").count() == 0
    assert not shard_dir("T3").exists()


# ---------------------------------------------------------------- run_task 接线


@pytest.fixture
def wired_pipeline(session, storage, monkeypatch):
    """把 run_task 需要的三个外部依赖接到测试会话上，再把入队换成记录器。"""
    monkeypatch.setattr(pipeline_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(
        pipeline_module,
        "select_engine",
        lambda meta, size_bytes, requested=None, graph_configured=False: "graph",
    )
    # I2：run_task 进分片分支前会真的调 is_graph_configured(session)——这个
    # 假 select_engine 不会替 run_task 挡住这一步，得让它在真实凭证表里
    # 查到能通过的一行，否则这里所有测试都会在 prepare_shards 之前被新增
    # 的早退检查拦下（见 GRAPH_TEST_SECRET_KEY 的说明）。
    monkeypatch.setattr(settings, "secret_key", GRAPH_TEST_SECRET_KEY)
    save_credentials(
        session,
        tenant_id="tid",
        client_id="cid",
        client_secret="secret",
        site_id="site-1",
        drive_path="staging",
    )
    enqueued: list[tuple[str, list[str]]] = []
    monkeypatch.setattr(
        pipeline_module,
        "enqueue_shards",
        lambda task_id, shard_ids: enqueued.append((task_id, list(shard_ids))),
    )
    return enqueued


def test_run_task_routes_to_sharded_path(session, deck, wired_pipeline, monkeypatch):
    """接线断言：删掉 run_task 里的 prepare_shards 调用，行不会建；
    删掉 enqueue_shards 调用，记录器是空的。两者任一都会让本测试变红。"""
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 2)
    task = _make_task(session, deck, task_id="T4")
    engine = _FakeEngine()
    asked = _install_engine(monkeypatch, engine, target=pipeline_module)

    pipeline_module.run_task("T4")

    assert task.shard_total == 4
    assert session.query(TaskShard).filter(TaskShard.task_id == "T4").count() == 4
    assert wired_pipeline == [("T4", [s.shard_id for s in _shards(session, "T4")])]
    # 分片路径不走单次转换
    assert engine.calls == []
    assert asked == []
    # 任务留在 converting，等汇总 job 落终态
    assert session.get(Task, "T4").status == "converting"


def test_run_task_keeps_single_conversion_when_not_sharding(
    session, deck, wired_pipeline, monkeypatch
):
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 100)
    monkeypatch.setattr(settings, "graph_max_shard_bytes", 500 * 1024 * 1024)
    task = _make_task(session, deck, task_id="T5")
    engine = _FakeEngine(pages=DECK_SLIDES)
    _install_engine(monkeypatch, engine, target=pipeline_module)

    pipeline_module.run_task("T5")

    assert wired_pipeline == []
    assert task.shard_total is None
    assert len(engine.calls) == 1
    assert session.get(Task, "T5").status == "done"


def test_run_task_fails_loudly_instead_of_falling_back(
    session, deck, wired_pipeline, monkeypatch
):
    """用户显式选了 Graph 而条件不满足时必须明确报错，
    绝不能偷偷改用 LibreOffice。"""
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", DECK_SLIDES)
    monkeypatch.setattr(settings, "graph_max_shard_bytes", 1)
    task = _make_task(session, deck, task_id="T6")
    engine = _FakeEngine()
    asked = _install_engine(monkeypatch, engine, target=pipeline_module)

    pipeline_module.run_task("T6")

    task = session.get(Task, "T6")
    assert task.status == "failed"
    assert task.error_code == "SHARD_TOO_LARGE"
    assert task.engine == "graph"  # 引擎没有被悄悄换掉
    assert engine.calls == []
    assert asked == []
    assert wired_pipeline == []


def test_run_task_fails_loudly_when_graph_not_configured(
    session, deck, wired_pipeline, monkeypatch
):
    """项目红线：用户显式选了 Graph，但 Azure 凭证根本没配置——必须明确
    报错，绝不能悄悄换成 LibreOffice 或占位引擎。

    刻意不用 _install_engine：这条测试要走真实的 get_engine("graph",
    session=...) 路径，验证 Task 8 从 GraphEngine.convert() 移出来的
    load_credentials 调用真的在 run_task 里被触发了，而不是只在
    test_engines_registry.py 里单独测过、run_task 那行接线却被删掉也没人
    发现。额外 patch httpx.Client 到会炸的桩，证明这条路径连 HTTP 都没碰到
    就已经失败——不是"发了请求、认证失败"，是根本没准备好就报错。
    """
    import app.services.engines.graph as graph_module

    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 100)
    monkeypatch.setattr(settings, "graph_max_shard_bytes", 500 * 1024 * 1024)
    monkeypatch.setattr(settings, "secret_key", None)  # 未配置

    def _boom(*a, **kw):
        raise AssertionError("不该发出任何 HTTP 请求——凭证检查必须先失败")

    monkeypatch.setattr(graph_module.httpx, "Client", _boom)

    task = _make_task(session, deck, task_id="T9")

    pipeline_module.run_task("T9")

    task = session.get(Task, "T9")
    assert task.status == "failed"
    assert task.error_code == "GRAPH_NOT_CONFIGURED"
    assert task.engine == "graph"  # 引擎没有被悄悄换掉
    assert wired_pipeline == []


def test_run_task_cleans_up_when_enqueue_fails(session, deck, wired_pipeline, monkeypatch):
    """入队失败（Redis 不可用）时，TaskShard 行已经 commit、最多 480MB 分片
    pptx 已经落盘。不清理的话它们是永久孤儿：merge_shards 永不执行，而
    purge_expired_outputs 只扫 outputs_dir，没有任何东西会清 shards_dir。"""
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 2)
    _make_task(session, deck, task_id="T7")

    def _boom(task_id: str, shard_ids: list[str]) -> None:
        raise ConnectionError("redis 连不上")

    monkeypatch.setattr(pipeline_module, "enqueue_shards", _boom)
    _install_engine(monkeypatch, _FakeEngine(), target=pipeline_module)

    pipeline_module.run_task("T7")

    task = session.get(Task, "T7")
    assert task.status == "failed"
    assert session.query(TaskShard).filter(TaskShard.task_id == "T7").count() == 0
    assert task.shard_total is None
    assert not shard_dir("T7").exists()


# --------------------------------------------- C1: requested_engine 必须真的传到 select_engine


def test_run_task_requested_graph_wins_even_when_auto_would_pick_libreoffice(
    session, deck, storage, monkeypatch
):
    """审查 Critical C1：`Task.requested_engine` 这个字段必须真的传到
    `select_engine` 的 `requested=` 形参。此前全仓所有 run_task 测试都靠
    `_force_placeholder_engine`/`wired_pipeline` 把 `select_engine` 整个换
    成忽略 `requested` 的常量 lambda，`requested_engine` 从未在任何断言里
    起过作用——如果 `pipeline.py` 里 `select_engine(meta, size_bytes,
    requested=task.requested_engine, ...)` 那个 `requested=` 被删掉，
    207 个测试一条不红。

    这里刻意不 patch select_engine，让真实实现跑；只 fake get_engine
    （避免真 HTTP/子进程）。deck 页数刻意设成超过自动路由阈值，让"如果
    requested 被漏传"和"正确传了"这两种情况给出不同结果——这才是有区分力
    的测试，不是恰好碰巧结果一样。
    """
    monkeypatch.setattr(pipeline_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    # conftest 的 _force_placeholder_engine 是 autouse 的，会把 select_engine
    # 整个换成忽略 requested 的常量 lambda——这里必须再 patch 回真实实现，
    # 否则测的就不是 requested= 这条接线本身了。
    monkeypatch.setattr(pipeline_module, "select_engine", real_select_engine)
    # 8 页的 deck 超过这个阈值——如果 requested 被漏传，auto 分支会判定
    # 页数超限，选出 libreoffice；正确传了的话，requested="graph" 直接
    # 短路掉 auto 判定，结果是 graph。
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 2)
    monkeypatch.setattr(
        pipeline_module,
        "enqueue_shards",
        lambda task_id, shard_ids: None,
    )
    _install_engine(monkeypatch, _FakeEngine(), target=pipeline_module)

    _make_fresh_task(session, deck, task_id="TC1A", requested_engine="graph")

    pipeline_module.run_task("TC1A")

    task = session.get(Task, "TC1A")
    assert task.engine == "graph"  # 没有被 auto 判定悄悄换成 libreoffice


def test_run_task_requested_libreoffice_wins_even_when_auto_would_pick_graph(
    session, deck, storage, monkeypatch
):
    """C1 的镜像场景：显式选 libreoffice，deck 又小又在阈值内、Graph 也
    配置好了——auto 分支如果被误调用会选 graph，但 requested= 必须让它
    连 auto 判定都不经过，直接短路成 libreoffice。"""
    monkeypatch.setattr(pipeline_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(pipeline_module, "select_engine", real_select_engine)
    monkeypatch.setattr(settings, "secret_key", GRAPH_TEST_SECRET_KEY)
    save_credentials(
        session,
        tenant_id="tid",
        client_id="cid",
        client_secret="secret",
        site_id="site-1",
        drive_path="staging",
    )
    engine = _FakeEngine(pages=DECK_SLIDES)
    asked = _install_engine(monkeypatch, engine, target=pipeline_module)

    _make_fresh_task(session, deck, task_id="TC1B", requested_engine="libreoffice")

    pipeline_module.run_task("TC1B")

    task = session.get(Task, "TC1B")
    assert task.engine == "libreoffice"
    assert asked == ["libreoffice"]  # 真的用假引擎转了一次，不是断言个寂寞
    assert task.status == "done"


def test_run_task_auto_selects_graph_when_configured_and_within_thresholds(
    session, deck, storage, monkeypatch
):
    """复审第二轮 Important：`graph_configured=` 这个传参本身此前无测试
    守护——上面两条 C1 测试都显式设了 requested_engine，`if requested:`
    直接短路，auto 分支（真正读取 graph_configured 的地方）从未被这两条
    测试真正跑到过。这里补上唯一会读到它的场景：requested_engine=None，
    走真实 select_engine 的 auto 分支，用小 deck（页数体积都在默认阈值
    内）+ 真实配置好的 Graph 凭证，断言 auto 判定真的选出了 "graph"——
    如果 `pipeline.py` 里 `graph_configured=graph_configured,` 那行传参
    被删掉（关键字专属参数，默认 False），auto 分支会因为拿到默认值
    `False` 而只会选 libreoffice，这里必须变红。"""
    monkeypatch.setattr(pipeline_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(pipeline_module, "select_engine", real_select_engine)
    monkeypatch.setattr(settings, "secret_key", GRAPH_TEST_SECRET_KEY)
    save_credentials(
        session,
        tenant_id="tid",
        client_id="cid",
        client_secret="secret",
        site_id="site-1",
        drive_path="staging",
    )
    # 页数体积都在默认阈值内（DECK_SLIDES=8 << graph_max_pages_per_shard=80，
    # deck 体积 << graph_max_shard_bytes=40MiB），不需要额外 monkeypatch 阈值。
    _install_engine(monkeypatch, _FakeEngine(), target=pipeline_module)

    _make_fresh_task(session, deck, task_id="TC1C", requested_engine=None)

    pipeline_module.run_task("TC1C")

    task = session.get(Task, "TC1C")
    assert task.engine == "graph"


# --------------------------------------------- I2: 分片前的凭证早退检查


def test_run_task_fails_before_splitting_when_graph_not_configured(
    session, deck, storage, monkeypatch
):
    """审查 Important I2：凭证检查必须在切片之前做，不能等 prepare_shards
    整轮切完（落盘 + commit TaskShard 行 + 入队）才在子 job 里第一次发现
    没配置。用一个会在被调用时立刻失败的 `prepare_shards` 桩断言它压根
    没被调用——如果早退检查被删掉或挪到 prepare_shards 之后，这里必须
    变红。"""
    monkeypatch.setattr(pipeline_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(retention_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(shard_module, "SessionLocal", lambda: session)
    monkeypatch.setattr(pipeline_module, "select_engine", real_select_engine)
    monkeypatch.setattr(settings, "graph_max_pages_per_shard", 2)  # 强制走分片分支
    monkeypatch.setattr(settings, "secret_key", None)  # 未配置

    def _prepare_shards_must_not_be_called(*a, **kw):
        raise AssertionError("凭证未配置时不该切片——早退检查必须先失败")

    monkeypatch.setattr(shard_module, "prepare_shards", _prepare_shards_must_not_be_called)

    _make_fresh_task(session, deck, task_id="TC-I2", requested_engine="graph")

    pipeline_module.run_task("TC-I2")

    task = session.get(Task, "TC-I2")
    assert task.status == "failed"
    assert task.error_code == "GRAPH_NOT_CONFIGURED"
    assert task.shard_total is None  # prepare_shards 真的没跑
    assert not shard_dir("TC-I2").exists()


def _shards(session, task_id: str) -> list[TaskShard]:
    return (
        session.query(TaskShard)
        .filter(TaskShard.task_id == task_id)
        .order_by(TaskShard.index)
        .all()
    )


# ---------------------------------------------------------------- 队列接线


def test_enqueue_shards_uses_dependency_with_allow_failure(monkeypatch):
    """allow_failure=True 是必须的：默认的 False 会让任一分片失败时汇总 job
    永远停在 DeferredJobRegistry 里不执行，任务卡死在 converting。"""
    from unittest.mock import Mock

    from rq.job import Dependency, Job

    import app.queue as queue_module
    from app.services.shard_pipeline import convert_shard as cs
    from app.services.shard_pipeline import merge_shards as ms

    class _FakeQueue:
        """只记录 enqueue 调用，不碰 Redis。返回真 Job 实例——Dependency
        会拒绝非 Job/str 的对象，用假对象测不到真实约束。"""

        def __init__(self) -> None:
            self.calls: list[tuple] = []

        def enqueue(self, func, *args, **kwargs):
            job = Job(id=f"job-{len(self.calls)}", connection=Mock())
            self.calls.append((func, args, kwargs, job))
            return job

    q = _FakeQueue()
    monkeypatch.setattr(queue_module, "get_queue", lambda: q)

    queue_module.enqueue_shards("T9", ["a", "b", "c"])

    assert len(q.calls) == 4
    for i, sid in enumerate(["a", "b", "c"]):
        func, args, kwargs, _job = q.calls[i]
        assert func is cs
        assert args == (sid,)
        assert kwargs["job_timeout"] > 0

    func, args, kwargs, _job = q.calls[3]
    assert func is ms
    assert args == ("T9",)
    dep = kwargs["depends_on"]
    assert isinstance(dep, Dependency)
    assert dep.allow_failure is True
    assert dep.dependencies == [c[3] for c in q.calls[:3]]
